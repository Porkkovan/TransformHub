"""
TransformHub – Senior Leadership Demo Deck
Talking points appear TWO ways:
  1. As a clearly styled yellow/amber NOTES STRIP at the bottom of every slide
  2. In the PowerPoint Notes pane (visible in Presenter View)
White slide backgrounds, indigo/cyan brand colours.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand palette ──────────────────────────────────────────────────────────
INDIGO   = RGBColor(0x4F, 0x46, 0xE5)
CYAN     = RGBColor(0x06, 0xB6, 0xD4)
GREEN    = RGBColor(0x16, 0xA3, 0x4A)
AMBER    = RGBColor(0xD9, 0x77, 0x06)
RED      = RGBColor(0xDC, 0x26, 0x26)
PURPLE   = RGBColor(0x7C, 0x3A, 0xED)
DARK     = RGBColor(0x1E, 0x29, 0x3B)
MID      = RGBColor(0x47, 0x55, 0x69)
LGRAY    = RGBColor(0xF1, 0xF5, 0xF9)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
TP_BG    = RGBColor(0xFF, 0xF7, 0xED)   # amber tint for talking points strip
TP_BORD  = RGBColor(0xF5, 0x9E, 0x0B)
INDBG    = RGBColor(0xEE, 0xF2, 0xFF)
CYANBG   = RGBColor(0xEC, 0xFE, 0xFF)
GREENBG  = RGBColor(0xF0, 0xFD, 0xF4)
AMBERBG  = RGBColor(0xFF, 0xFB, 0xEB)
REDBG    = RGBColor(0xFE, 0xF2, 0xF2)

W = Inches(13.33)
H = Inches(7.5)

# Vertical zones (inches)
CONTENT_TOP    = 1.15   # content starts below header
CONTENT_BOTTOM = 5.10   # content ends here
TP_TOP         = 5.20   # talking points strip starts
TP_H           = 1.95   # height of TP strip → ends at 7.15
FOOTER_TOP     = 7.18

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ══════════════════════════════════════════════════════════════════════════════
# Core drawing helpers
# ══════════════════════════════════════════════════════════════════════════════

def new_slide():
    sl = prs.slides.add_slide(BLANK)
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = WHITE
    return sl

def rect(sl, x, y, w, h, fill=None, line=None, lw=Pt(0.75)):
    sh = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line; sh.line.width = lw
    else:
        sh.line.fill.background()
    return sh

def oval(sl, x, y, w, h, fill):
    sh = sl.shapes.add_shape(9, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh

def tb(sl, text, x, y, w, h, sz=11, bold=False, color=DARK,
       align=PP_ALIGN.LEFT, italic=False):
    box = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return box

def multiline(sl, lines, x, y, w, h, sz=10, color=DARK, spacing_inches=0.26):
    for i, line in enumerate(lines):
        tb(sl, line, x, y + i * spacing_inches, w, spacing_inches + 0.05,
           sz=sz, color=color)

# ══════════════════════════════════════════════════════════════════════════════
# Slide furniture
# ══════════════════════════════════════════════════════════════════════════════

def header(sl, title, sub=None, accent=INDIGO):
    rect(sl, 0, 0, 13.33, 0.08, fill=accent)
    tb(sl, title, 0.4, 0.15, 10, 0.58, sz=22, bold=True, color=DARK)
    if sub:
        tb(sl, sub, 0.4, 0.73, 12.5, 0.32, sz=10, color=MID, italic=True)

def footer(sl):
    rect(sl, 0, FOOTER_TOP, 13.33, 0.28, fill=LGRAY)
    tb(sl, "TransformHub  ·  Senior Leadership Demo  ·  March 2026  ·  http://localhost:3000",
       0.3, FOOTER_TOP + 0.04, 12.7, 0.2, sz=7.5, color=MID)

def add_notes(sl, points):
    """Write talking points to the PPT notes pane."""
    notes_tf = sl.notes_slide.notes_text_frame
    notes_tf.clear()
    p = notes_tf.paragraphs[0]
    r = p.add_run()
    r.text = "TALKING POINTS"
    r.font.bold = True; r.font.size = Pt(12)
    for pt in points:
        para = notes_tf.add_paragraph()
        run = para.add_run()
        run.text = f"•  {pt}"
        run.font.size = Pt(11)

def talking_points(sl, points, accent=AMBER):
    """
    Amber strip at bottom of every slide showing talking points.
    Clearly labelled, never overlaps content.
    """
    # Background strip
    rect(sl, 0, TP_TOP, 13.33, TP_H, fill=TP_BG, line=TP_BORD, lw=Pt(0))
    # Left colour tab
    rect(sl, 0, TP_TOP, 0.22, TP_H, fill=TP_BORD)
    # Label
    tb(sl, "💬  TALKING POINTS", 0.3, TP_TOP + 0.08, 2.5, 0.3,
       sz=9, bold=True, color=TP_BORD)
    # Points in two columns
    mid = len(points) // 2 + len(points) % 2
    col1 = points[:mid]
    col2 = points[mid:]
    xs = [0.3, 6.75]
    for ci, col in enumerate([col1, col2]):
        x = xs[ci]
        for ri, pt in enumerate(col):
            y = TP_TOP + 0.42 + ri * 0.46
            oval(sl, x + 0.02, y + 0.08, 0.10, 0.10, fill=TP_BORD)
            tb(sl, pt, x + 0.20, y, 6.2, 0.44, sz=9, color=DARK)
    # Notes pane also
    add_notes(sl, points)

def metric(sl, val, label, x, y, w=1.7, h=0.9,
           vcol=INDIGO, bg=INDBG, border=INDIGO):
    rect(sl, x, y, w, h, fill=bg, line=border, lw=Pt(0.6))
    tb(sl, str(val), x, y + 0.07, w, 0.46, sz=26, bold=True,
       color=vcol, align=PP_ALIGN.CENTER)
    tb(sl, label, x, y + 0.54, w, 0.3, sz=8.5, color=MID,
       align=PP_ALIGN.CENTER)

def card(sl, x, y, w, h, bg=LGRAY, border=INDIGO, lw=Pt(0.6)):
    return rect(sl, x, y, w, h, fill=bg, line=border, lw=lw)

def step_badge(sl, num, x, y, color=INDIGO, sz=0.34):
    c = sl.shapes.add_shape(9, Inches(x), Inches(y), Inches(sz), Inches(sz))
    c.fill.solid(); c.fill.fore_color.rgb = color; c.line.fill.background()
    tf = c.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = str(num); r.font.size = Pt(12); r.font.bold = True
    r.font.color.rgb = WHITE

def acc_bar(sl, label, pct, x, y, w=4.8, bar_col=GREEN):
    tb(sl, label, x, y, 3.2, 0.22, sz=9.5, color=DARK)
    tb(sl, f"{pct}%", x + 3.3, y, 0.7, 0.22, sz=9.5, bold=True,
       color=bar_col, align=PP_ALIGN.RIGHT)
    rect(sl, x, y + 0.25, w, 0.14, fill=LGRAY)
    fw = max(0.05, w * pct / 100)
    rect(sl, x, y + 0.25, fw, 0.14, fill=bar_col)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1  Title
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 0.22, 7.5, fill=INDIGO)   # left accent band

tb(sl, "TransformHub", 0.5, 0.6, 9.5, 1.0,
   sz=48, bold=True, color=INDIGO)
tb(sl, "Enterprise Digital Transformation Platform  ·  Senior Leadership Demo  ·  March 2026",
   0.5, 1.62, 11.0, 0.38, sz=12, color=MID)
tb(sl, "Powered by 18 AI Agents  ·  Next.js 15 + FastAPI + LangGraph  ·  PostgreSQL 18 + pgvector",
   0.5, 2.05, 11.0, 0.32, sz=10.5, color=MID, italic=True)

stats = [("3","Enterprise\nDemo Orgs"), ("8","Core\nModules"),
         ("18","AI\nAgents"), ("78","Capabilities\nMapped"), ("245+","Functionalities\nDiscovered")]
for i,(v,l) in enumerate(stats):
    metric(sl, v, l, 0.5 + i*2.0, 2.6, 1.7, 1.05)

rect(sl, 0.5, 3.85, 8.2, 0.9, fill=INDBG, line=INDIGO, lw=Pt(0.6))
tb(sl, "AI-powered transformation — from discovery to delivery — fully auditable, benchmark-grounded and transparent.",
   0.7, 3.93, 7.9, 0.7, sz=11, italic=True, color=INDIGO)

# Credentials block
rect(sl, 9.3, 0.55, 3.8, 4.5, fill=INDBG, line=INDIGO, lw=Pt(0.75))
rect(sl, 9.3, 0.55, 3.8, 0.42, fill=INDIGO)
tb(sl, "DEMO ACCESS", 9.5, 0.62, 3.4, 0.28, sz=11, bold=True, color=WHITE)
creds = [
    ("Pre-Filled Demo","demo@transformhub.ai","demo1234",INDBG),
    ("Live Entry Demo","live@transformhub.ai","live1234",CYANBG),
    ("Admin","admin@transformhub.ai","demo1234",LGRAY),
]
cy=1.08
for label,email,pwd,bg in creds:
    rect(sl,9.45,cy,3.5,0.92,fill=bg,line=INDIGO,lw=Pt(0.4))
    tb(sl,label,9.6,cy+0.05,3.2,0.26,sz=9,bold=True,color=INDIGO)
    tb(sl,f"✉  {email}",9.6,cy+0.30,3.2,0.24,sz=9,color=DARK)
    tb(sl,f"🔑  {pwd}",9.6,cy+0.54,3.2,0.24,sz=9,color=MID)
    cy+=1.02
tb(sl,"🌐  http://localhost:3000",9.5,cy+0.08,3.4,0.3,sz=10,bold=True,color=CYAN)

footer(sl)
talking_points(sl,[
    "TransformHub consolidates 6+ separate enterprise tools into one AI-native platform",
    "18 LangGraph agents cover the full transformation lifecycle — discovery through delivery",
    "Multi-tenant: US Bank, Telstra Health and ING Bank are all isolated and independently seeded",
    "78 capabilities mapped with 85% average confidence — evidence-based, not AI guesswork",
    "245 functionalities with full source triangulation (3–5 evidence sources each)",
    "Accuracy dashboard shows live composite score (~68%) — this rises as more agents are run",
])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2  Agenda
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl,"Today's Agenda","25-minute structured walkthrough across all platform modules")

agenda=[
    (INDIGO,"1","Platform Overview","3 min","Architecture, 8 modules, 18 AI agents pipeline"),
    (CYAN,  "2","Discovery Module","4 min","Multi-pass AI · confidence scoring · source triangulation"),
    (GREEN, "3","Value Stream Mapping","4 min","L1/L2/L3 hierarchy · flow efficiency · Mermaid diagrams"),
    (PURPLE,"4","Future State Vision","4 min","Benchmark-grounded projections · automation mix"),
    (RED,   "5","Risk & Compliance","3 min","Risk gates · SHA-256 audit trail · regulatory frameworks"),
    (AMBER, "6","Accuracy Dashboard","3 min","Per-module scores · composite · live action plan"),
    (CYAN,  "7","Live Entry Demo","4 min","Create org · run agents · watch data populate live"),
    (GREEN, "8","Q & A","5 min","Open questions · next steps · roadmap"),
]
for i,(col,num,title,dur,desc) in enumerate(agenda):
    c=i%2; r=i//2
    xb=0.35+c*6.55; yb=CONTENT_TOP+r*0.93
    card(sl,xb,yb,6.3,0.83,bg=LGRAY,border=col,lw=Pt(0.7))
    step_badge(sl,num,xb+0.1,yb+0.22,color=col,sz=0.36)
    tb(sl,title,xb+0.60,yb+0.06,3.5,0.36,sz=12,bold=True,color=DARK)
    tb(sl,f"⏱ {dur}",xb+4.2,yb+0.06,1.8,0.3,sz=9,bold=True,color=col,align=PP_ALIGN.RIGHT)
    tb(sl,desc,xb+0.60,yb+0.44,5.5,0.3,sz=9.5,color=MID)

footer(sl)
talking_points(sl,[
    "Keep each section tight to the time allocation — the 25-min format is designed for exec attention spans",
    "Start with pre-filled demo (Steps 1–6) to show polish — add live entry at the end if time allows",
    "Steps 2–6 each have accuracy badges on the page header — point to them as you navigate",
    "Q&A (Step 8) is where you ask: 'Which of your products would you like us to map first?'",
    "If they're engaged, offer to run the Discovery agent live against their actual repo during Q&A",
])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3  Demo Mode Selection
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl,"Choose Your Demo Path","Select based on audience type and time available")

for ci,(col,bg,icon,mode,tag,items) in enumerate([
    (INDIGO,INDBG,"🎯","Pre-Filled Demo","Recommended for Leadership",[
        "Login: demo@transformhub.ai  /  demo1234",
        "3 enterprise orgs: US Bank · Telstra Health · ING Bank",
        "78 capabilities, 245 functionalities, full VSM seeded",
        "Accuracy scores visible immediately (~65–72% composite)",
        "21 context docs across all 5 RAG categories — all INDEXED",
        "234 roadmap items with RICE scores and approval status",
        "Estimated duration: 15–20 minutes",
    ]),
    (CYAN,CYANBG,"⚡","Live Entry Demo","Best for Technical Audiences",[
        "Login: live@transformhub.ai  /  live1234",
        "Create a new org — blank slate, no pre-seeded data",
        "Paste any GitHub repo URL or public website URL",
        "Run Discovery → review → run VSM → check Accuracy",
        "Watch accuracy scores rise from 0% to live values in real-time",
        "Requires OPENAI_API_KEY in agent-service/.env",
        "Estimated duration: 20–30 minutes",
    ]),
]):
    xb=0.35+ci*6.55
    card(sl,xb,CONTENT_TOP,6.3,3.8,bg=bg,border=col)
    rect(sl,xb,CONTENT_TOP,6.3,0.44,fill=col)
    tb(sl,f"{icon}  {mode}",xb+0.18,CONTENT_TOP+0.08,5.9,0.3,sz=13,bold=True,color=WHITE)
    tb(sl,tag,xb+0.18,CONTENT_TOP+0.56,5.9,0.28,sz=10,bold=True,color=col)
    for ri,item in enumerate(items):
        tb(sl,f"✓  {item}",xb+0.18,CONTENT_TOP+0.9+ri*0.42,5.9,0.38,sz=10,color=DARK)

footer(sl)
talking_points(sl,[
    "For Board / C-suite: always use Pre-Filled — the visual richness of 78 capabilities + 234 roadmap items is the story",
    "For IT / Architecture teams: Live Entry is more impressive — they understand what real-time agent execution means",
    "Pre-Filled shows 3 different industry verticals (Banking, Healthcare, Payments) — switch orgs to show multi-tenancy",
    "Live Entry works best with a public GitHub repo — stripe/stripe-node is ideal (clear API structure, well-tested)",
    "Both modes use the same agent pipeline — the difference is only whether data is pre-seeded or generated live",
])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4  Platform Architecture
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl,"Platform Architecture","Full-stack enterprise AI platform — three tiers, production-grade")

layers=[
    (INDIGO,"⚛  Frontend  :3000","Next.js 15 App Router  ·  TypeScript  ·  Tailwind v4  ·  Dark glassmorphism UI  ·  Real-time streaming output"),
    (CYAN,  "🐍  Agent Service  :8000","FastAPI  ·  LangGraph  ·  Python 3.13  ·  18 AI Agents  ·  BM25 + pgvector RAG  ·  Agent memory + feedback loops"),
    (GREEN, "🗄  Data Layer  :5432","PostgreSQL 18  ·  pgvector extension  ·  Prisma ORM  ·  SHA-256 chained audit trail  ·  Multi-tenant row isolation"),
]
for i,(col,title,desc) in enumerate(layers):
    yb=CONTENT_TOP+i*1.22
    card(sl,0.35,yb,12.6,1.08,bg=LGRAY,border=col)
    rect(sl,0.35,yb,0.22,1.08,fill=col)
    tb(sl,title,0.7,yb+0.12,5.5,0.38,sz=13,bold=True,color=col)
    tb(sl,desc,0.7,yb+0.58,11.8,0.38,sz=10,color=DARK)
    if i<2:
        tb(sl,"▼",6.4,yb+1.1,0.6,0.25,sz=12,color=col,align=PP_ALIGN.CENTER)

# Stats row
stats=[("18","LangGraph Agents",INDIGO),("50+","API Routes",CYAN),
       ("8","Core Modules",GREEN),("5","RAG Categories",AMBER),("3","Tenants",PURPLE)]
for i,(v,l,c) in enumerate(stats):
    metric(sl,v,l,0.35+i*2.55,CONTENT_TOP+3.75,2.3,0.9,vcol=c,bg=LGRAY,border=c)

footer(sl)
talking_points(sl,[
    "Three-tier architecture means each layer can be scaled independently — agent service scales horizontally for heavy workloads",
    "LangGraph checkpointing means if an agent fails mid-run, it resumes from the last checkpoint — no lost work",
    "Redis is optional — the platform runs fully on in-memory fallback, which is what you see in the demo",
    "pgvector enables the hybrid RAG: BM25 for keyword recall + cosine similarity for semantic recall — best of both",
    "50+ API routes cover the full CRUD surface for all 8 modules plus agent execution, accuracy, and context management",
    "Prisma ORM gives type-safe DB access — no raw SQL, no injection risk, strong migration tooling",
])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5  8 Core Modules
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl,"8 Core Modules","End-to-end transformation lifecycle — Discovery through Delivery")

modules=[
    (INDIGO,"🔍","Discovery","Multi-pass AI analysis\n3 human review gates\nConfidence + source triangulation\n~85% accuracy score"),
    (GREEN, "🗺","Value Stream\nMapping","L1 / L2 / L3 hierarchy\nMermaid diagrams auto-generated\n100% VSM coverage\n~72% accuracy score"),
    (PURPLE,"🚀","Future State\nVision","Benchmark-grounded projections\n3-band metrics: cons/exp/opt\nAutomation mix breakdown\n~70% accuracy score"),
    (RED,   "🛡","Risk &\nCompliance","Risk gate: score ≥8.0 blocks\nSHA-256 chained audit trail\nSOX · FDIC · HIPAA · Basel III"),
    (AMBER, "🗓","Product\nRoadmap","RICE scoring algorithm\nApproval workflow (PENDING→APPROVED)\n234 items, Q1–Q4 2026"),
    (CYAN,  "🔧","Product\nWorkbench","Readiness scoring 0–10\nArchitecture views (functional/tech)\nInline editing of all entities"),
    (GREEN, "📚","Context\nHub","5 RAG doc categories\nBM25 + vector hybrid search\nURL fetch + file upload"),
    (AMBER, "📊","Accuracy\nDashboard","Per-module composite score\nLive action plan (4–8 steps)\nBadges on every page header"),
]
for i,(col,icon,title,desc) in enumerate(modules):
    c=i%4; r=i//4
    xb=0.35+c*3.24; yb=CONTENT_TOP+r*1.95
    card(sl,xb,yb,3.08,1.78,bg=LGRAY,border=col)
    rect(sl,xb,yb,3.08,0.1,fill=col)
    tb(sl,icon,xb+0.12,yb+0.18,0.5,0.44,sz=18)
    tb(sl,title,xb+0.65,yb+0.22,2.3,0.44,sz=11,bold=True,color=col)
    tb(sl,desc,xb+0.12,yb+0.7,2.85,0.98,sz=8.5,color=MID)

footer(sl)
talking_points(sl,[
    "Each module feeds the next: Discovery → VSM → Future State → Risk → Roadmap is a complete AI-driven workflow",
    "Accuracy badges appear on EVERY module page header — click any badge to jump to the Accuracy Dashboard",
    "Context Hub is the brain behind all agents — docs uploaded there are injected into every agent's RAG context",
    "Discovery is the entry point — without running Discovery first, other agents have limited data to work with",
    "Product Workbench and Roadmap are the output modules — this is where leadership makes go/no-go decisions",
    "Risk & Compliance has a hard gate at score ≥8.0 — no transformation proceeds until risk is mitigated",
])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6  Discovery
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl,"Step 1  ·  Discovery Module","Multi-pass AI analysis with human review gates — ~85% accuracy",accent=INDIGO)

passes=[
    (INDIGO,"Pass 1","Digital Products identified","Human reviews product list → Edit / Delete inline → Approve to continue"),
    (CYAN,  "Pass 2","Capabilities mapped","Constrained to approved products only → human reviews capability tree"),
    (GREEN, "Pass 3","Functionalities extracted","Confidence score + evidence sources per item → human final approval"),
]
for i,(col,pnum,title,desc) in enumerate(passes):
    yb=CONTENT_TOP+i*1.18
    card(sl,0.35,yb,8.0,1.0,bg=LGRAY,border=col)
    rect(sl,0.35,yb,0.22,1.0,fill=col)
    step_badge(sl,str(i+1),0.65,yb+0.31,color=col,sz=0.34)
    tb(sl,pnum,1.12,yb+0.08,2.0,0.34,sz=11,bold=True,color=col)
    tb(sl,title,1.12,yb+0.44,3.0,0.3,sz=11,color=DARK)
    tb(sl,desc,4.3,yb+0.08,3.9,0.82,sz=9.5,color=MID,italic=True)
    if i<2: tb(sl,"▼",4.1,yb+1.02,0.5,0.2,sz=11,color=col,align=PP_ALIGN.CENTER)

# Source types
rect(sl,8.55,CONTENT_TOP,4.45,3.52,fill=INDBG,line=INDIGO,lw=Pt(0.6))
rect(sl,8.55,CONTENT_TOP,4.45,0.36,fill=INDIGO)
tb(sl,"EVIDENCE SOURCES",8.72,CONTENT_TOP+0.06,4.1,0.24,sz=9,bold=True,color=WHITE)
sources=[("GitHub Structure","Code architecture + file tree"),
         ("OpenAPI Spec","Endpoint + schema definitions"),
         ("DB Schema","Table + relationship inference"),
         ("GitHub Tests","Behaviour documentation in tests"),
         ("URL Analysis","Public-facing product pages"),
         ("Context Documents","Uploaded BRDs and process maps"),
         ("Integration Data","Connected system metadata"),
         ("Q&A / Questionnaire","Analyst-provided context"),]
for ri,(src,desc) in enumerate(sources):
    y=CONTENT_TOP+0.45+ri*0.37
    oval(sl,8.7,y+0.1,0.1,0.1,fill=INDIGO)
    tb(sl,src,8.88,y,2.0,0.26,sz=9,bold=True,color=DARK)
    tb(sl,desc,10.95,y,1.9,0.26,sz=8.5,color=MID)

# Stats
stats=[("9","Products"),("26","Capabilities",GREEN),("79","Functionalities"),
       ("85%","Avg Confidence",GREEN),("100%","Triangulated",GREEN)]
xst=0.35
for item in stats:
    v,l=item[0],item[1]; c=item[2] if len(item)>2 else INDIGO
    metric(sl,v,l,xst,CONTENT_TOP+3.6,1.58,0.85,vcol=c,bg=LGRAY,border=c)
    xst+=1.65

footer(sl)
talking_points(sl,[
    "CLICK: Discovery → switch to Drilldown view → click 'LoanFlow Digital' → show 9 capabilities expanding in the tree",
    "Point to the confidence badge on 'AI Underwriting' capability — 87% confidence from 4 evidence sources (GitHub + OpenAPI + DB + docs)",
    "CLICK: switch to Product Catalog view — this is the exec-friendly summary: 3 products, 26 capabilities, 79 functionalities at a glance",
    "Inline editing: hover over any capability name → pencil icon appears → edit live → changes persist immediately",
    "Multi-pass design prevents hallucination cascade — AI cannot invent capabilities for products that haven't been approved",
    "Source triangulation: 100% of seeded capabilities have 3+ evidence sources — show the source distribution breakdown in Accuracy page",
])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7  VSM
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl,"Step 2  ·  Value Stream Mapping","Three-level hierarchy — 100% capability coverage — Mermaid visualisation",accent=GREEN)

levels=[
    (GREEN, "L1","Segment Level","Cross-product flow efficiency view · Business segment selector · Comparative flow efficiency bars"),
    (CYAN,  "L2","Product\nCapabilities","Process Time / Wait Time / Lead Time / Flow Efficiency per capability · Mermaid diagram rendered"),
    (INDIGO,"L3","Functionality\nSteps","Step-by-step timing · Classification: Value-Adding / Bottleneck / Waste / Waiting · Drill from L2"),
]
for i,(col,lvl,title,desc) in enumerate(levels):
    yb=CONTENT_TOP+i*1.2
    card(sl,0.35,yb,7.8,1.05,bg=LGRAY,border=col)
    rect(sl,0.35,yb,0.22,1.05,fill=col)
    tb(sl,lvl,0.65,yb+0.27,0.7,0.5,sz=18,bold=True,color=col,align=PP_ALIGN.CENTER)
    tb(sl,title,1.5,yb+0.1,2.5,0.5,sz=13,bold=True,color=DARK)
    tb(sl,desc,4.2,yb+0.1,3.8,0.82,sz=10,color=MID)
    if i<2: tb(sl,"▼",4.0,yb+1.08,0.5,0.2,sz=11,color=col,align=PP_ALIGN.CENTER)

# Key metrics explanation
rect(sl,8.45,CONTENT_TOP,4.55,3.52,fill=GREENBG,line=GREEN,lw=Pt(0.6))
rect(sl,8.45,CONTENT_TOP,4.55,0.36,fill=GREEN)
tb(sl,"KEY VSM METRICS",8.62,CONTENT_TOP+0.06,4.2,0.24,sz=9,bold=True,color=WHITE)
vsm_defs=[
    ("PT  Process Time","Time actively working on the step (value-added time)"),
    ("WT  Wait Time","Queue / idle time before the next step begins"),
    ("LT  Lead Time","Total elapsed time = PT + WT"),
    ("FE  Flow Efficiency","PT ÷ LT × 100 — industry avg ~38%"),
    ("Bottleneck","Step where WT > PT and blocks downstream flow"),
    ("Waste","Steps with no value-add — candidates for elimination"),
]
for ri,(term,def_) in enumerate(vsm_defs):
    y=CONTENT_TOP+0.45+ri*0.49
    tb(sl,term,8.62,y,2.1,0.26,sz=9,bold=True,color=GREEN)
    tb(sl,def_,8.62,y+0.26,4.25,0.24,sz=8.5,color=MID)

stats=[("78","Capabilities with VSM"),("100%","Coverage",GREEN),
       ("78","Mermaid Diagrams"),("~35%","Avg Flow Efficiency",AMBER),("60%+","Avg Waste",RED)]
xst=0.35
for item in stats:
    v,l=item[0],item[1]; c=item[2] if len(item)>2 else INDIGO
    metric(sl,v,l,xst,CONTENT_TOP+3.6,1.58,0.85,vcol=c,bg=LGRAY,border=c)
    xst+=1.65

footer(sl)
talking_points(sl,[
    "CLICK: VSM → set Level = L1 → show flow efficiency bars across all 3 products — LoanFlow is the bottleneck at ~28%",
    "CLICK: L2 → select 'AI Underwriting' → point to the Mermaid diagram auto-rendered below the metrics table",
    "Average Flow Efficiency ~35% means 65% of lead time is WASTE — this is the transformation opportunity",
    "CLICK: Classification panel (right side of L2) → show Bottleneck steps highlighted in orange, Waste in red",
    "100% VSM coverage = all 78 capabilities have metrics — previous session had only 1 per product (3 total)",
    "CLICK: Capability Comparison Chart button → shows side-by-side FE bars for all capabilities in the segment",
])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8  Future State
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl,"Step 3  ·  Future State Vision","Benchmark-grounded projections with automation mix and 3-band metrics",accent=PURPLE)

features=[
    (PURPLE,"Automation Mix Breakdown","Per-product breakdown across: RPA · AI/ML · Agent-Based · Conversational Analytics\nShows technology investment split as stacked bar chart"),
    (AMBER, "Projected Metrics — 3 Bands","Conservative / Expected / Optimistic projections grounded in uploaded VSM benchmark docs\n'Benchmark-grounded' badge confirms projections are data-backed, not AI estimates"),
    (CYAN,  "Capability Modernisation Roadmap","RICE scoring per capability: Reach × Impact × Confidence ÷ Effort\nBusiness impact · Complexity rating · Recommended tech stack · Auto-feeds Product Roadmap agent"),
]
for i,(col,title,desc) in enumerate(features):
    yb=CONTENT_TOP+i*1.2
    card(sl,0.35,yb,12.6,1.05,bg=LGRAY,border=col)
    rect(sl,0.35,yb,0.22,1.05,fill=col)
    tb(sl,title,0.68,yb+0.1,4.0,0.38,sz=13,bold=True,color=col)
    tb(sl,desc,0.68,yb+0.52,11.7,0.45,sz=10,color=DARK)
    if i<2: tb(sl,"▼",6.4,yb+1.08,0.5,0.2,sz=11,color=col,align=PP_ALIGN.CENTER)

rect(sl,0.35,CONTENT_TOP+3.65,12.6,0.78,fill=AMBERBG,line=AMBER,lw=Pt(0.6))
tb(sl,"💡  What 'Benchmark-Grounded' means:",0.55,CONTENT_TOP+3.72,5.0,0.3,sz=10,bold=True,color=AMBER)
tb(sl,"When VSM_BENCHMARKS or TRANSFORMATION_CASE_STUDIES documents are uploaded to Context Hub, the Future State agent reads them and derives "
   "projections from real industry data rather than generic multipliers. A 'Benchmark-grounded' badge appears in the header confirming this.",
   0.55,CONTENT_TOP+4.0,12.2,0.36,sz=9.5,color=DARK)

stats=[("9","Products with Vision"),("3","Metric Bands"),("✓","Benchmark Docs",GREEN),("234","Roadmap Items"),("~70%","FS Accuracy",GREEN)]
xst=0.35
for item in stats:
    v,l=item[0],item[1]; c=item[2] if len(item)>2 else PURPLE
    metric(sl,v,l,xst,CONTENT_TOP+4.5,2.3,0.85,vcol=c,bg=LGRAY,border=c)
    xst+=2.42

footer(sl)
talking_points(sl,[
    "CLICK: Future State → note the 'Benchmark-grounded' badge in the header — hover it to see which docs are grounding the projections",
    "CLICK: Automation Mix chart → show LoanFlow Digital: 30% RPA, 40% AI/ML, 20% Agent-Based — ask 'does this match your current roadmap thinking?'",
    "CLICK: Expand any product → Projected Metrics table → point to Conservative/Expected/Optimistic bands with specific FE targets",
    "The Conservative band assumes 20% of benchmark improvement, Expected = 50%, Optimistic = 80% — conservative is the planning number",
    "CLICK: Capability Modernisation list → sort by RICE score → top item is the highest ROI transformation candidate",
    "Future State output auto-saves as an AGENT_OUTPUT context doc — the Roadmap agent reads it to generate roadmap items",
])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9  Risk & Compliance
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl,"Step 4  ·  Risk & Compliance","Hard risk gates · SHA-256 tamper-proof audit trail · regulatory framework mapping",accent=RED)

sev=[
    (RED,  "CRITICAL  ≥ 8.0","BLOCKS transformation gate. Mitigation must be completed and approved before any capability moves to production."),
    (AMBER,"HIGH  6.0 – 7.9","Transformation allowed with an approved mitigation plan. Plan tracked as mandatory pre-condition in roadmap."),
    (INDIGO,"MEDIUM  4.0 – 5.9","Advisory. Mitigation recommended but does not block the transformation gate."),
    (GREEN,"LOW  < 4.0","Informational. Logged in audit trail. No action required before transformation proceeds."),
]
for i,(col,label,desc) in enumerate(sev):
    yb=CONTENT_TOP+i*0.88
    card(sl,0.35,yb,8.0,0.78,bg=LGRAY,border=col)
    rect(sl,0.35,yb,0.22,0.78,fill=col)
    tb(sl,label,0.68,yb+0.1,2.2,0.3,sz=11,bold=True,color=col)
    tb(sl,desc,0.68,yb+0.44,7.5,0.3,sz=9.5,color=DARK)

# Frameworks
rect(sl,8.55,CONTENT_TOP,4.45,3.52,fill=REDBG,line=RED,lw=Pt(0.6))
rect(sl,8.55,CONTENT_TOP,4.45,0.36,fill=RED)
tb(sl,"REGULATORY FRAMEWORKS",8.72,CONTENT_TOP+0.06,4.1,0.24,sz=9,bold=True,color=WHITE)
orgs_fw=[
    ("US Bank","FINRA · SEC · SOX · FDIC · BSA/AML"),
    ("Telstra Health","HL7 FHIR · HIPAA · GDPR · TGA"),
    ("ING Bank","Basel III · MiFID II · SOLVENCY II · DORA"),
    ("SHA-256 Audit","Each entry hashes prev — tamper-proof chain"),
    ("Risk Agent","Auto-generates mitigation plan per risk item"),
    ("Gate Logic","Score ≥ 8.0 → transitionBlocked = TRUE in DB"),
]
for ri,(org,fws) in enumerate(orgs_fw):
    y=CONTENT_TOP+0.44+ri*0.49
    tb(sl,org,8.72,y,1.9,0.26,sz=9,bold=True,color=RED)
    tb(sl,fws,8.72,y+0.26,4.25,0.24,sz=8.5,color=MID)

# Audit trail explanation
rect(sl,0.35,CONTENT_TOP+3.55,8.0,0.88,fill=REDBG,line=RED,lw=Pt(0.5))
tb(sl,"🔒  SHA-256 Chained Audit Trail",0.55,CONTENT_TOP+3.62,4.0,0.3,sz=10,bold=True,color=RED)
tb(sl,"Each log entry stores: Action · Actor · Payload · PayloadHash(previous+current). "
   "If any entry is altered, the hash chain breaks — immediately detectable. Required for SOX/FDIC/GDPR compliance.",
   0.55,CONTENT_TOP+3.92,7.7,0.44,sz=9.5,color=DARK)

stats=[("12","Risk Assessments"),("18","Compliance Maps"),("3","Frameworks / Org"),("4","Audit Entries / Org"),("SHA-256","Chain Type")]
xst=0.35
for v,l in stats:
    c=RED if "Risk" in l else INDIGO
    metric(sl,v,l,xst,CONTENT_TOP+4.5,2.3,0.85,vcol=c,bg=LGRAY,border=c)
    xst+=2.42

footer(sl)
talking_points(sl,[
    "CLICK: Risk & Compliance → point to the CRITICAL risk item — show the 'Transition BLOCKED' badge on the capability card",
    "Explain the gate: until this risk score drops below 8.0, no transformation agent will move that capability forward",
    "CLICK: Compliance Mappings tab → show FDIC + SOX requirements mapped to specific capabilities with evidence links",
    "CLICK: Audit Log tab → scroll through SHA-256 chained entries — point to 'previousHash' field linking entries together",
    "Ask the audience: 'For your most regulated product, what frameworks would you need mapped?' — write it down for follow-up",
    "The Risk agent takes ~60 seconds to run — offer to run it live during Q&A to show real-time risk identification",
])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10  Accuracy Dashboard
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl,"Step 5  ·  Accuracy Dashboard","Per-module AI output quality with composite scoring and live action plan",accent=AMBER)

# Bars (left column)
bars=[
    ("Discovery",                 85, GREEN),
    ("Lean VSM",                  72, GREEN),
    ("Future State Vision",       70, GREEN),
    ("Knowledge Base / RAG",      70, GREEN),
    ("Product Transformation",    55, AMBER),
    ("Risk & Compliance",         35, AMBER),
    ("Architecture",              40, AMBER),
]
tb(sl,"MODULE SCORES",0.35,CONTENT_TOP+0.05,4.0,0.28,sz=9,bold=True,color=MID)
for i,(label,pct,col) in enumerate(bars):
    acc_bar(sl,label,pct,0.35,CONTENT_TOP+0.40+i*0.52,w=5.0,bar_col=col)

# Composite
rect(sl,0.35,CONTENT_TOP+4.10,5.3,0.6,fill=INDBG,line=INDIGO,lw=Pt(0.75))
tb(sl,"COMPOSITE SCORE  (weighted average)",0.55,CONTENT_TOP+4.18,4.5,0.26,sz=9,bold=True,color=INDIGO)
acc_bar(sl,"Discovery×20% + VSM×18% + FutureState×15% + KB×15% + Risk×12% + PT×12% + Arch×8%",
        68,0.35,CONTENT_TOP+4.5,w=5.3,bar_col=INDIGO)

# Scoring legend
rect(sl,5.9,CONTENT_TOP,3.0,2.0,fill=LGRAY,line=INDIGO,lw=Pt(0.5))
rect(sl,5.9,CONTENT_TOP,3.0,0.36,fill=INDIGO)
tb(sl,"SCORE LEGEND",6.08,CONTENT_TOP+0.07,2.7,0.24,sz=9,bold=True,color=WHITE)
legend=[(GREEN,"80–100%","Excellent"),(GREEN,"65–79%","Good"),
        (AMBER,"45–64%","Fair"),(RED,"0–44%","Needs Work")]
for ri,(col,rng,lbl) in enumerate(legend):
    y=CONTENT_TOP+0.46+ri*0.38
    rect(sl,6.08,y,0.2,0.22,fill=col)
    tb(sl,f"{rng}  —  {lbl}",6.36,y,2.4,0.28,sz=10,color=DARK)

# Badge explanation
rect(sl,5.9,CONTENT_TOP+2.15,3.0,1.38,fill=AMBERBG,line=AMBER,lw=Pt(0.5))
rect(sl,5.9,CONTENT_TOP+2.15,3.0,0.35,fill=AMBER)
tb(sl,"ACCURACY BADGES",6.08,CONTENT_TOP+2.21,2.7,0.24,sz=9,bold=True,color=WHITE)
tb(sl,"Every module page header shows a live accuracy badge (green/amber/red). "
   "Click any badge to jump directly to this Accuracy Dashboard.",
   6.08,CONTENT_TOP+2.58,2.78,0.82,sz=9,color=DARK)

# Action plan
rect(sl,9.15,CONTENT_TOP,4.1,4.68,fill=AMBERBG,line=AMBER,lw=Pt(0.75))
rect(sl,9.15,CONTENT_TOP,4.1,0.38,fill=AMBER)
tb(sl,"⚡  TOP IMPROVEMENT ACTIONS",9.33,CONTENT_TOP+0.07,3.8,0.26,sz=9,bold=True,color=WHITE)
actions=[
    ("HIGH","Run Risk Agent → push Risk 35% → 80%+",
     "Risk agent covers all 26 capabilities vs current 4"),
    ("HIGH","Upload VSM Benchmarks to Context Hub",
     "Unlocks benchmark-grounded Future State projections"),
    ("HIGH","Add OpenAPI + DB schema as enrichment sources",
     "Each source adds 15–20% confidence per entity"),
    ("MED", "Run Architecture agent",
     "Architecture module currently at 0% — one run → ~60%"),
    ("MED", "Upload ARCHITECTURE_STANDARDS docs",
     "KB category missing → adds 8pts to KB composite"),
]
cy=CONTENT_TOP+0.50
for pri,title,sub in actions:
    col=RED if pri=="HIGH" else AMBER
    rect(sl,9.25,cy,0.42,0.24,fill=col)
    tb(sl,pri,9.25,cy,0.42,0.24,sz=7.5,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    tb(sl,title,9.74,cy,3.4,0.26,sz=9,bold=True,color=DARK)
    tb(sl,sub,9.74,cy+0.28,3.4,0.24,sz=8.5,color=MID,italic=True)
    cy+=0.75

footer(sl)
talking_points(sl,[
    "CLICK: Accuracy from sidebar → point to the composite score first: ~68% 'Good' — then explain why each module scored as it did",
    "Discovery 85% = high confidence avg (85%) + 100% source triangulation + 7 unique evidence types",
    "Risk 35% = only 4 assessments mapped vs 26 capabilities — 'run the Risk agent now and watch this jump to 80%+'",
    "CLICK: Expand the Discovery module card → show source distribution breakdown (GitHub, OpenAPI, DB, tests etc.)",
    "CLICK: Scroll to Action Plan → show the 5 prioritised improvement steps — each one links to the relevant module",
    "The badges on every page header make accuracy visible at all times — analysts never have to wonder about data quality",
])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11  Live Entry Demo
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl,"Step 6  ·  Live Entry Demo (Optional)","Create a new org and watch AI agents populate data from scratch",accent=CYAN)

steps=[
    (CYAN,  "1","Login as Live User","live@transformhub.ai  /  live1234","Dashboard empty — no pre-seeded data, blank slate"),
    (CYAN,  "2","Create Organisation","Organizations → + New Organization","Enter company name, industry type, business segments"),
    (CYAN,  "3","Run Discovery — Pass 1","Discovery → Paste GitHub URL → Analyze","Recommended: github.com/stripe/stripe-node (public, well-structured)"),
    (CYAN,  "4","Review Pass 1","Products stream in → Edit inline → Approve","Human review gate — edit product names to match your terminology"),
    (GREEN, "5","Run Pass 2 + 3","Capabilities then functionalities auto-run","Each entity gets confidence score + evidence sources"),
    (GREEN, "6","Run Lean VSM Agent","VSM page → Run Lean VSM","~30s — flow metrics + Mermaid diagrams generated for all caps"),
    (INDIGO,"7","Check Accuracy Page","Accuracy from sidebar","Watch composite rise from 0% → ~45% after first agent cycle"),
]
for i,(col,num,action,detail,tip) in enumerate(steps):
    c=i%2; r=i//2
    xb=0.35+c*6.55; yb=CONTENT_TOP+r*0.9
    if i==6: xb=0.35  # last item full width
    w=6.3 if i<6 else 12.8
    card(sl,xb,yb,w,0.78,bg=LGRAY,border=col)
    step_badge(sl,num,xb+0.1,yb+0.21,color=col,sz=0.34)
    tb(sl,action,xb+0.6,yb+0.06,3.2,0.3,sz=11,bold=True,color=col)
    tb(sl,detail,xb+0.6,yb+0.38,3.0,0.28,sz=9.5,color=DARK)
    tb(sl,f"💡 {tip}",xb+(w/2),yb+0.06,w/2-0.2,0.62,sz=9,color=MID,italic=True)

# Good repos
rect(sl,0.35,CONTENT_TOP+3.62,12.8,0.78,fill=CYANBG,line=CYAN,lw=Pt(0.6))
tb(sl,"RECOMMENDED DEMO REPOS  (all public):",0.55,CONTENT_TOP+3.70,4.0,0.26,sz=9,bold=True,color=CYAN)
repos=[("github.com/stripe/stripe-node","Payments — clear API, well-documented, multi-capability"),
       ("github.com/plaid/plaid-python","Fintech — financial data, strong compliance signals"),
       ("github.com/twilio/twilio-python","Communications — multi-product, good for discovery")]
for ri,(repo,why) in enumerate(repos):
    tb(sl,f"→  {repo}",0.55+ri*4.3,CONTENT_TOP+4.02,3.5,0.26,sz=9.5,bold=True,color=DARK)
    tb(sl,why,0.55+ri*4.3,CONTENT_TOP+4.26,3.5,0.22,sz=8.5,color=MID)

footer(sl)
talking_points(sl,[
    "Login as live@transformhub.ai — point out the empty dashboard: 'No data — we're going to fill this with AI agents in real-time'",
    "Create the org first: use a company name the audience knows or respects — makes the output immediately relatable",
    "Paste stripe/stripe-node GitHub URL — it's public, well-structured, and Discovery returns clean results in ~30 seconds",
    "During Pass 1 streaming — narrate what's happening: 'The agent is reading the repo structure, identifying product boundaries...'",
    "After Pass 1 approval, note the green confidence badges: 'These are evidence-based, not guesses — each product is backed by 3–4 source types'",
    "Show the Accuracy page at the very end: composite jumps from 0% to ~45% after one Discovery + VSM cycle — real-time quality scoring",
])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12  Business Value
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl,"Business Value","Why TransformHub matters for enterprise transformation programmes")

pillars=[
    (INDIGO,"⏱","Speed","Hours, Not Months",
     "Multi-pass AI discovery maps an enterprise product portfolio in under 10 minutes.\n"
     "Traditional baselining takes 6–12 months of manual workshops and analyst time."),
    (GREEN, "🎯","Accuracy","Evidence-Based Confidence",
     "8 evidence source triangulation. 78–97% confidence per entity.\n"
     "Human-in-the-loop review gates at every step — no black-box AI."),
    (AMBER, "💰","Cost Reduction","One Platform, Not Six",
     "Replaces: process mining, risk management, roadmap planning,\n"
     "compliance tracking, architecture review, and knowledge management tools."),
    (CYAN,  "🔒","Compliance","Audit-Ready by Design",
     "SHA-256 chained immutable audit trail.\n"
     "SOX, FDIC, GDPR, HIPAA, Basel III, MiFID II framework coverage built in."),
]
for i,(col,icon,title,subtitle,body) in enumerate(pillars):
    c=i%2; r=i//2
    xb=0.35+c*6.55; yb=CONTENT_TOP+r*2.0
    card(sl,xb,yb,6.3,1.82,bg=LGRAY,border=col)
    rect(sl,xb,yb,6.3,0.1,fill=col)
    tb(sl,icon,xb+0.15,yb+0.18,0.55,0.45,sz=22)
    tb(sl,title,xb+0.78,yb+0.18,2.5,0.38,sz=14,bold=True,color=col)
    tb(sl,subtitle,xb+0.78,yb+0.58,5.3,0.3,sz=10.5,bold=True,color=DARK)
    tb(sl,body,xb+0.15,yb+0.96,5.95,0.74,sz=9.5,color=MID)

footer(sl)
talking_points(sl,[
    "SPEED: A typical transformation programme spends 30–40% of its budget just on current-state baselining — TransformHub eliminates this",
    "ACCURACY: Every entity has a confidence score and source list — when an auditor asks 'how do you know this capability exists?', show them the 4 evidence sources",
    "COST: A conservative estimate: replacing 6 tools at £50k/year each = £300k/year saved, plus the analyst headcount previously doing manual baselining",
    "COMPLIANCE: SHA-256 chaining means the audit trail is legally defensible — every action is timestamped, hashed, and linked to the previous entry",
    "The Human-in-the-Loop design is deliberate: AI proposes, humans approve. This is critical for regulated industries where AI alone cannot make decisions",
    "All 4 pillars are visible in the platform — offer to show any one of them live in Q&A (accuracy → Discovery page, compliance → Risk page, etc.)",
])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13  Pre-Demo Checklist
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl,"Pre-Demo Checklist","Run this 15 minutes before the meeting — not 1 minute before")

# Services checklist
card(sl,0.35,CONTENT_TOP,5.85,3.75,bg=LGRAY,border=GREEN)
rect(sl,0.35,CONTENT_TOP,5.85,0.36,fill=GREEN)
tb(sl,"✓  SERVICES RUNNING",0.55,CONTENT_TOP+0.07,5.4,0.24,sz=10,bold=True,color=WHITE)
service_checks=[
    ("Next.js frontend","http://localhost:3000  →  redirects to /login"),
    ("FastAPI backend","curl localhost:8000/api/v1/health  →  {status: healthy}"),
    ("PostgreSQL","Postgres.app in menu bar shows green dot"),
    ("Demo data seeded","3 orgs in switcher: US Bank · Telstra Health · ING Bank"),
    ("Accuracy page loads","Composite score visible (~68%) for US Bank"),
]
for ri,(item,detail) in enumerate(service_checks):
    y=CONTENT_TOP+0.48+ri*0.62
    rect(sl,0.5,y+0.05,0.22,0.22,fill=WHITE,line=GREEN,lw=Pt(0.75))
    tb(sl,item,0.82,y,2.3,0.28,sz=10,bold=True,color=DARK)
    tb(sl,detail,0.82,y+0.30,4.9,0.26,sz=8.5,color=MID,italic=True)

# Browser checklist
card(sl,6.45,CONTENT_TOP,6.5,3.75,bg=LGRAY,border=INDIGO)
rect(sl,6.45,CONTENT_TOP,6.5,0.36,fill=INDIGO)
tb(sl,"✓  BROWSER SETUP",6.65,CONTENT_TOP+0.07,6.1,0.24,sz=10,bold=True,color=WHITE)
browser_checks=[
    ("Open incognito window","Prevents cached session from a different account causing issues"),
    ("Zoom to 90%","Ensures all table columns are fully visible without horizontal scroll"),
    ("Pre-logged in as demo","demo@transformhub.ai / demo1234  —  US Bank loaded"),
    ("KPI cards populated","Dashboard shows: 9 products, 26 capabilities, 79 functionalities"),
    ("This guide open","Open demo-guide.html in a separate browser tab for reference"),
]
for ri,(item,detail) in enumerate(browser_checks):
    y=CONTENT_TOP+0.48+ri*0.62
    rect(sl,6.6,y+0.05,0.22,0.22,fill=WHITE,line=INDIGO,lw=Pt(0.75))
    tb(sl,item,6.92,y,2.8,0.28,sz=10,bold=True,color=DARK)
    tb(sl,detail,6.92,y+0.30,5.8,0.26,sz=8.5,color=MID,italic=True)

# Start commands
rect(sl,0.35,CONTENT_TOP+3.88,6.35,0.88,fill=RGBColor(0x0F,0x17,0x2A))
tb(sl,"# Terminal 1 — Backend",0.55,CONTENT_TOP+3.96,5.9,0.22,sz=8.5,color=MID,italic=True)
tb(sl,"cd TransformHub/agent-service && source venv/bin/activate",
   0.55,CONTENT_TOP+4.18,5.9,0.22,sz=9,color=CYAN)
tb(sl,"uvicorn app.main:app --reload --port 8000",
   0.55,CONTENT_TOP+4.40,5.9,0.22,sz=9,color=CYAN)
tb(sl,"# Terminal 2 — Frontend   cd TransformHub/nextjs-app && npm run dev",
   0.55,CONTENT_TOP+4.62,5.9,0.22,sz=8.5,color=WHITE)

# Quick fixes
rect(sl,6.95,CONTENT_TOP+3.88,5.95,0.88,fill=AMBERBG,line=AMBER,lw=Pt(0.6))
tb(sl,"⚠  QUICK FIXES",7.15,CONTENT_TOP+3.93,5.5,0.24,sz=9,bold=True,color=AMBER)
fixes=[("Wrong org:","localStorage.removeItem('currentOrgId') in console"),
       ("Port 3000 busy:","kill $(lsof -ti:3000)  then restart Next.js"),
       ("No data:","cd nextjs-app && npx tsx prisma/seed.ts")]
for ri,(prob,fix) in enumerate(fixes):
    y=CONTENT_TOP+4.22+ri*0.24
    tb(sl,prob,7.15,y,1.5,0.22,sz=9,bold=True,color=AMBER)
    tb(sl,fix,8.7,y,3.95,0.22,sz=9,color=DARK)

footer(sl)
talking_points(sl,[
    "The most common demo failure: wrong org auto-loads. Fix: localStorage.removeItem('currentOrgId') in browser console, then refresh",
    "Port 3000 is shared with Product Coach (another project) — if it's running, kill it first: kill $(lsof -ti:3000)",
    "Run the health check BEFORE the meeting: curl localhost:8000/api/v1/health — if database shows disconnected, restart Postgres.app",
    "Redis unavailable in the health response is NORMAL — the platform uses in-memory fallback, no impact on demo",
    "Have a backup slide ready: if the live demo fails mid-run, switch to the pre-filled US Bank org immediately",
    "Pre-position the browser on the Dashboard with US Bank loaded — start the meeting already in the app, not on the login page",
])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14  Closing
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl,0,0,0.22,7.5,fill=INDIGO)

tb(sl,"Ready to Transform.",0.5,0.65,10.5,0.95,sz=42,bold=True,color=INDIGO)
tb(sl,"AI-powered enterprise transformation — from discovery to delivery.",
   0.5,1.62,10.5,0.42,sz=14,color=MID,italic=True)

summary=[("3","Enterprise\nDemo Orgs",INDIGO),("8","Core\nModules",CYAN),
         ("18","AI\nAgents",GREEN),("78","Capabilities\nMapped",INDIGO),
         ("245+","Functionalities",INDIGO),("~68%","Composite\nAccuracy",GREEN)]
for i,(v,l,c) in enumerate(summary):
    metric(sl,v,l,0.5+i*2.12,2.35,1.88,1.0,vcol=c,bg=LGRAY,border=c)

rect(sl,0.5,3.6,12.3,1.55,fill=INDBG,line=INDIGO,lw=Pt(0.75))
tb(sl,"🌐  http://localhost:3000",0.7,3.70,5.0,0.34,sz=13,bold=True,color=INDIGO)
for i,(label,cred) in enumerate([
    ("Pre-Filled Demo:","demo@transformhub.ai  /  demo1234   — Use for leadership demos — 3 enterprise orgs, full data"),
    ("Live Entry Demo:", "live@transformhub.ai  /  live1234   — Use for technical walkthroughs — blank org, run agents live"),
    ("Admin:","admin@transformhub.ai  /  demo1234"),
]):
    tb(sl,label,0.7,3.73+(i+1)*0.33,2.5,0.3,sz=9.5,bold=True,color=MID)
    tb(sl,cred,3.3,3.73+(i+1)*0.33,9.3,0.3,sz=9.5,color=DARK)

badges=["Next.js 15 + TypeScript","FastAPI + LangGraph","PostgreSQL 18 + pgvector",
        "18 AI Agents","BM25 + Vector RAG","SHA-256 Audit Trail","Multi-Tenant"]
x=0.5
for b in badges:
    w=len(b)*0.092+0.28
    rect(sl,x,5.42,w,0.38,fill=INDBG,line=INDIGO,lw=Pt(0.4))
    tb(sl,b,x+0.1,5.47,w-0.15,0.26,sz=9,color=INDIGO)
    x+=w+0.1

footer(sl)
talking_points(sl,[
    "Close with: 'Which of your products would you want us to map first in Discovery?'",
    "If budget conversations come up: ROI is hours-to-discovery vs months of analyst workshops — the payback is in the first sprint",
    "Next steps: provide their GitHub repo URL or website URL and we can have a Discovery output ready in 24 hours",
    "The platform is live at localhost:3000 — leave the browser open so they can explore after the meeting",
    "Offer to re-seed with their company name and industry so the demo feels personalised next time",
])

# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════
out="/Users/125066/projects/TransformHub/docs/TransformHub_Demo_Deck.pptx"
prs.save(out)
print(f"✅  Saved: {out}")
print(f"   {len(prs.slides)} slides")
