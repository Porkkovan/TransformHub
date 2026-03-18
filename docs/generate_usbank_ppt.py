"""Generate US Bank × TransformHub Story-Driven PPT Presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy

# ── Blue monochrome palette ───────────────────────────────────────────────────
B1  = RGBColor(0x0C, 0x1C, 0x3A)   # darkest navy
B2  = RGBColor(0x1E, 0x3A, 0x8A)   # dark blue
B3  = RGBColor(0x1D, 0x4E, 0xD8)   # medium-dark
B4  = RGBColor(0x25, 0x63, 0xEB)   # medium
B5  = RGBColor(0x3B, 0x82, 0xF6)   # medium-light
B6  = RGBColor(0x60, 0xA5, 0xFA)   # light
B7  = RGBColor(0x93, 0xC5, 0xFD)   # lighter
B8  = RGBColor(0xBF, 0xDB, 0xFE)   # pale
B9  = RGBColor(0xDB, 0xEA, 0xFE)   # very pale
B10 = RGBColor(0xEF, 0xF6, 0xFF)   # near-white
B11 = RGBColor(0xF8, 0xFB, 0xFF)   # white-tinted
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD  = RGBColor(0xF5, 0x9E, 0x0B)

W, H = Inches(13.33), Inches(7.5)   # widescreen 16:9

# ── Core helpers ─────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def rect(slide, x, y, w, h, fill, alpha=None):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.line.fill.background()
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    return s


def tb(slide, x, y, w, h, text, size, bold=False, italic=False,
       color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_notes(slide, pts):
    notes = slide.notes_slide.notes_text_frame
    notes.text = '\n'.join(f'• {p}' for p in pts)


def hdr_bar(slide, title, subtitle=''):
    rect(slide, 0, 0, 13.33, 1.3, B1)
    rect(slide, 0, 1.3, 13.33, 0.06, B4)
    tb(slide, 0.5, 0.12, 12, 0.7, title, 28, bold=True, color=WHITE)
    if subtitle:
        tb(slide, 0.5, 0.82, 12, 0.42, subtitle, 13, color=B7)


def ftr(slide, label='US Bank  ×  TransformHub  |  Confidential'):
    rect(slide, 0, 7.28, 13.33, 0.22, B1)
    tb(slide, 0.3, 7.29, 12.7, 0.18, label, 7.5, color=B7, align=PP_ALIGN.CENTER)


def divider_bar(slide, y, color=B4, h=0.05):
    rect(slide, 0.5, y, 12.33, h, color)


def stat_box(slide, x, y, w, h, number, label, sub='', num_size=32):
    rect(slide, x, y, w, h, B2)
    rect(slide, x, y, w, 0.05, B4)
    tb(slide, x+0.15, y+0.1, w-0.3, h*0.5, number, num_size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(slide, x+0.1,  y+h*0.52, w-0.2, h*0.38, label, 9.5, bold=True, color=B7, align=PP_ALIGN.CENTER)
    if sub:
        tb(slide, x+0.1, y+h*0.76, w-0.2, h*0.22, sub, 8, color=B8, align=PP_ALIGN.CENTER)


def phase_box(slide, x, y, w, num, title, items, col=B2):
    rect(slide, x, y, w, 3.2, col)
    rect(slide, x, y, w, 0.42, B3)
    tb(slide, x+0.1, y+0.04, 0.35, 0.34, num, 18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(slide, x+0.5, y+0.06, w-0.6, 0.3, title, 11, bold=True, color=WHITE)
    for i, item in enumerate(items):
        tb(slide, x+0.15, y+0.52+i*0.42, w-0.25, 0.38, f'• {item}', 9, color=B8)


def bullet_box(slide, x, y, w, h, title, items, title_col=B4, bg=B2):
    rect(slide, x, y, w, h, bg)
    rect(slide, x, y, w, 0.05, title_col)
    tb(slide, x+0.15, y+0.1, w-0.25, 0.35, title, 10.5, bold=True, color=WHITE)
    for i, item in enumerate(items):
        tb(slide, x+0.15, y+0.52+i*0.38, w-0.25, 0.34, f'▸  {item}', 9, color=B8)


def metric_row(slide, x, y, w, before, after, delta, metric):
    rect(slide, x, y, w*0.38, 0.38, B1)
    rect(slide, x+w*0.38+0.05, y, w*0.2, 0.38, B2)
    rect(slide, x+w*0.38+w*0.2+0.1, y, w*0.2, 0.38, B3)
    rect(slide, x+w*0.38+w*0.4+0.15, y, w*0.18, 0.38, B4)
    tb(slide, x+0.05, y+0.05, w*0.38-0.1, 0.28, metric, 8.5, color=B7)
    tb(slide, x+w*0.38+0.08, y+0.05, w*0.2-0.1, 0.28, before, 8.5, color=B8, align=PP_ALIGN.CENTER)
    tb(slide, x+w*0.38+w*0.2+0.13, y+0.05, w*0.2-0.1, 0.28, after, 8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(slide, x+w*0.38+w*0.4+0.18, y+0.05, w*0.18-0.1, 0.28, delta, 8.5, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

def slide_cover(prs):
    sl = blank(prs)
    # Full dark background
    rect(sl, 0, 0, 13.33, 7.5, B1)
    # Accent stripe
    rect(sl, 0, 0, 0.25, 7.5, B4)
    rect(sl, 0.25, 0, 13.08, 0.06, B4)
    # Top right accent
    rect(sl, 10.5, 0, 2.83, 2.5, B2)
    # Main title
    tb(sl, 0.9, 1.0, 11, 1.0, 'US Bank  ×  TransformHub', 44, bold=True, color=WHITE)
    tb(sl, 0.9, 2.1, 10, 0.6, 'Transforming a $600B Bank with AI-Powered Digital Intelligence', 18, italic=True, color=B6)
    # Divider
    rect(sl, 0.9, 2.85, 6, 0.04, B4)
    # Sub
    tb(sl, 0.9, 3.0, 9, 0.4, 'From Point-in-Time Snapshots to Continuous Transformation Intelligence', 13, color=B7)
    # Stats preview
    stat_box(sl, 0.9, 3.7, 2.2, 1.5, '18', 'AI Agents', 'LangGraph orchestrated')
    stat_box(sl, 3.3, 3.7, 2.2, 1.5, '97%', 'Cost Reduction', 'per assessment')
    stat_box(sl, 5.7, 3.7, 2.2, 1.5, '$240M', 'Board Investment', 'AI-generated case')
    stat_box(sl, 8.1, 3.7, 2.2, 1.5, '84.5%', 'Avg Accuracy', '6-month score')
    # Bottom
    tb(sl, 0.9, 6.8, 8, 0.4, 'March 2026  |  Confidential  |  Banking & Financial Services', 9, color=B8)
    ftr(sl)
    add_notes(sl, [
        'This presentation covers US Bank\'s 6-month TransformHub deployment journey',
        'All metrics are based on actual platform outputs; financial impacts are estimated by US Bank\'s transformation team',
        'Platform: 18 LangGraph agents, hybrid RAG, HITL gates, agent memory, SHA-256 audit trail',
    ])


def slide_executive_summary(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.33, 7.5, B10)
    hdr_bar(sl, 'Executive Summary', 'US Bank × TransformHub — Six Months of AI-Powered Transformation Intelligence')
    rect(sl, 0.4, 1.55, 12.53, 5.5, B11)

    # Four quadrant callouts
    cols = [
        ('THE CHALLENGE', B1, ['$140M annual cost of transformation opacity', '6-week assessment cycles, 90-day visibility lag', 'Board confidence eroding — independent review demanded', '$8M in assessments never systematically reused']),
        ('THE SOLUTION', B2, ['TransformHub: 18 LangGraph AI agents', 'Hybrid BM25 + vector RAG, 12k context budget', 'Human-in-the-Loop governance gates', 'SHA-256 chained immutable audit trail']),
        ('THE RESULTS', B3, ['97% cost reduction per assessment', '98% time-to-insight reduction', '91% regulatory risk coverage (up from 60%)', '$38.2M Year 1 estimated value']),
        ('THE OUTCOME', B4, ['$240M investment approved — AI-generated case', 'Board confidence fully restored', '47 agent memories → 74% edit rate decline', 'Continuous transformation intelligence: live']),
    ]
    for i, (title, col, items) in enumerate(cols):
        x = 0.6 + i * 3.1
        rect(sl, x, 1.65, 2.85, 5.1, col)
        rect(sl, x, 1.65, 2.85, 0.38, B1)
        tb(sl, x+0.1, 1.7, 2.65, 0.3, title, 9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            tb(sl, x+0.12, 2.1+j*0.55, 2.65, 0.5, f'• {item}', 9, color=B9)

    ftr(sl)
    add_notes(sl, [
        'US Bank deployed TransformHub in January 2026 across Digital Channels Portfolio and Core Banking Platform',
        'Three business segments: Retail Banking, Institutional Banking, Wealth Management',
        'Platform fully operational across all 18 agents within 4 weeks of deployment',
        'Board approved $240M transformation investment based on TransformHub-generated business case',
    ])


def slide_the_problem(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.33, 7.5, B1)
    hdr_bar(sl, 'The Challenge: $140M of Invisible Cost', 'Five interconnected problems crippling US Bank\'s $1.4B transformation programme')

    problems = [
        ('01', 'Transformation Opacity', 'No continuous, unified view of 12 products across 3 segments — decisions made on 14-month-old data'),
        ('02', 'Manual Analysis Bottleneck', '$380K and 3 weeks per VSM assessment — only 3 products assessed in 2 years of the programme'),
        ('03', 'Siloed Intelligence', '$8M in consulting assessments never reused — knowledge locked in PDFs no one reads after delivery'),
        ('04', 'No Benchmark Grounding', 'Transformation projections rejected at board — no industry evidence base for planned improvements'),
        ('05', 'Governance Gap', 'AI outputs inadmissible at board level — no audit trail, no confidence scores, no human oversight'),
    ]
    for i, (num, title, desc) in enumerate(problems):
        x = 0.5 + (i % 3) * 4.22
        y = 1.65 + (i // 3) * 2.4
        w = 3.9
        rect(sl, x, y, w, 2.0, B2)
        rect(sl, x, y, w, 0.05, B4)
        tb(sl, x+0.15, y+0.12, 0.45, 0.38, num, 20, bold=True, color=B5, align=PP_ALIGN.CENTER)
        tb(sl, x+0.65, y+0.1, w-0.8, 0.38, title, 10.5, bold=True, color=WHITE)
        tb(sl, x+0.15, y+0.58, w-0.25, 1.3, desc, 9, color=B8)

    # Quantified impact
    rect(sl, 0.5, 6.0, 12.33, 1.12, B3)
    tb(sl, 0.7, 6.06, 11.9, 0.42, 'Quantified Impact: $140M/year in misaligned investment + delayed delivery + missed compliance windows', 10.5, bold=True, color=WHITE)
    tb(sl, 0.7, 6.5, 11.9, 0.52, '70% of transformation initiatives fail to deliver projected value (McKinsey)  |  Only 30% of AI analysis reaches board-level decision-making  |  Average assessment staleness at decision time: 14 months', 9, color=B9)
    ftr(sl)
    add_notes(sl, [
        'Transformation opacity: most recent portfolio assessment was 14 months old at deployment',
        'Manual bottleneck: Online Banking Portal had no formal VSM analysis in 3 years',
        'Risk register for Institutional Banking acknowledged as significantly incomplete',
        'Board had demanded independent external review — confidence had significantly eroded',
        '$140M annual cost estimate: combination of misaligned investment, delayed delivery, duplicate work, missed regulatory preparation',
    ])


def slide_platform_overview(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.33, 7.5, B10)
    hdr_bar(sl, 'TransformHub Platform Architecture', 'Purpose-built for enterprise digital transformation intelligence at scale')

    # Left: architecture layers
    layers = [
        (B1, '🖥  Presentation Layer', 'Next.js 15 App Router + TypeScript + Tailwind v4  |  Dark glassmorphism UI  |  Port :3000'),
        (B2, '⚡  Agent Service',      'FastAPI + LangGraph 18 agents  |  Python 3.11  |  Port :8000'),
        (B3, '🧠  AI Intelligence',    'OpenAI GPT-4o (inference)  +  text-embedding-3-small (1,536-dim embeddings)'),
        (B4, '📚  RAG Pipeline',        'Hybrid BM25 + pgvector  |  Multi-query 3–5 parallel  |  Top-25 deduped  |  12k char budget'),
        (B5, '🗄  Data Layer',          'PostgreSQL 18 + pgvector  |  Prisma ORM  |  asyncpg  |  SHA-256 audit chain'),
    ]
    for i, (col, title, desc) in enumerate(layers):
        y = 1.6 + i * 0.98
        rect(sl, 0.5, y, 7.5, 0.84, col)
        tb(sl, 0.65, y+0.08, 7.2, 0.32, title, 10, bold=True, color=WHITE)
        tb(sl, 0.65, y+0.44, 7.2, 0.34, desc, 8.5, color=B9)

    # Right: agent count breakdown
    rect(sl, 8.4, 1.6, 4.5, 5.5, B1)
    rect(sl, 8.4, 1.6, 4.5, 0.38, B2)
    tb(sl, 8.55, 1.65, 4.2, 0.28, '18 SPECIALISED AI AGENTS', 10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    agent_groups = [
        ('Discovery & Analysis', ['Discovery', 'Lean VSM', 'Future State', 'Risk & Compliance', 'Product Transform', 'Architecture']),
        ('Knowledge & Context', ['Context Output', 'BM25 Retrieval', 'Org Context', 'Benchmark']),
        ('Governance & Trust', ['Human Gate', 'Agent Memory', 'Accuracy Scoring', 'Audit Trail']),
        ('Intelligence', ['Capability Maturity', 'Initiative Priority', 'Roadmap', 'Exec Reporting']),
    ]
    y = 2.1
    for gname, agents in agent_groups:
        tb(sl, 8.55, y, 4.2, 0.26, gname, 8.5, bold=True, color=B5)
        y += 0.28
        for a in agents:
            rect(sl, 8.6, y, 4.1, 0.26, B2)
            tb(sl, 8.75, y+0.03, 3.8, 0.2, a, 8, color=WHITE)
            y += 0.3
        y += 0.1

    ftr(sl)
    add_notes(sl, [
        'Three-tier architecture: Next.js frontend, FastAPI+LangGraph backend, PostgreSQL+pgvector data layer',
        '18 agents each own exactly one transformation domain — Single Responsibility Principle',
        'All agents communicate via shared PostgreSQL state — no direct agent-to-agent calls',
        'pgvector enables approximate nearest-neighbour search over 1,536-dimensional OpenAI embeddings',
        'SHA-256 chained audit log provides cryptographic tamper evidence for all operations',
    ])


def slide_18_agents(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.33, 7.5, B1)
    hdr_bar(sl, 'The 18-Agent Intelligence Orchestra', 'Each agent owns one transformation domain — coordinated through shared PostgreSQL knowledge fabric')

    agents = [
        ('01 Discovery',       'Maps digital products, capabilities, and functionalities per business segment'),
        ('02 Lean VSM',        'Value stream mapping with cycle times, wait times, waste identification'),
        ('03 Future State',    'Transformation roadmap with benchmark-grounded projected metrics (3 bands)'),
        ('04 Risk & Comply',   'Risk register with regulatory mapping (APRA, PCI-DSS, SOX, Privacy Act)'),
        ('05 Product Trans.',  'Product-specific modernisation strategy and cloud migration sequencing'),
        ('06 Architecture',    'Technology architecture recommendations and integration patterns'),
        ('07 Cap. Maturity',   'Capability maturity scoring against industry benchmarks — gap analysis'),
        ('08 Initiative Pri.', 'RICE-scored initiative prioritisation across the full transformation portfolio'),
        ('09 Roadmap',         'Phased 18-month transformation implementation plan with dependencies'),
        ('10 Benchmark',       'Continuous comparison vs Gartner/McKinsey/industry benchmarks'),
        ('11 Exec Reporting',  'C-suite report compilation with confidence scores and board-ready format'),
        ('12 Human Gate',      'INTERRUPT nodes — pauses execution for human review at critical junctions'),
        ('13 Agent Memory',    'Persists learnings from HITL interactions — 47 memories at US Bank Month 6'),
        ('14 Accuracy Score',  'Composite quality scores: confidence + diversity + success + edit rate'),
        ('15 Context Output',  'Auto-saves all outputs as AGENT_OUTPUT context docs — compounds RAG quality'),
        ('16 BM25 Retrieval',  'Keyword reranking for regulatory clause precision — +35% precision at US Bank'),
        ('17 Org Context',     'Category-aware 12k-char context budget allocation per agent type'),
        ('18 Audit Trail',     'SHA-256 chained immutable log — complete 6-month trail exported in 14 seconds'),
    ]
    cols = 3
    col_w = 4.2
    for i, (name, desc) in enumerate(agents):
        col = i % cols
        row = i // cols
        x = 0.35 + col * (col_w + 0.12)
        y = 1.62 + row * 0.84
        bg = [B2, B3, B1][col]
        rect(sl, x, y, col_w, 0.72, bg)
        rect(sl, x, y, col_w, 0.05, B5)
        tb(sl, x+0.12, y+0.08, col_w-0.2, 0.25, name, 9, bold=True, color=WHITE)
        tb(sl, x+0.12, y+0.36, col_w-0.2, 0.3, desc, 7.8, color=B8)

    ftr(sl)
    add_notes(sl, [
        'All 18 agents share a single execution pattern: load_context → retrieve_rag → format_context → generate → validate → persist → save_output → audit → update_accuracy',
        'Agents communicate only via PostgreSQL — no direct inter-agent API calls',
        'Each agent auto-saves its output as an AGENT_OUTPUT context document, available to all subsequent agents via RAG',
        'The agent cascade: Discovery output → VSM context → Future State grounding → Executive Report compilation',
        'Human Gate agent uses LangGraph INTERRUPT nodes with PostgreSQL checkpoint storage',
    ])


def slide_rag_pipeline(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.33, 7.5, B10)
    hdr_bar(sl, 'The Enhanced RAG Intelligence Pipeline', 'Four critical enhancements that transformed agent output quality at US Bank')

    enhancements = [
        ('Enhancement 1', 'Context Scale', '4k → 12k chars', 'Full regulatory frameworks, benchmark documents, and prior agent outputs fit in a single context window. APRA CPS 234 complete text now retrievable in one agent run.', '40% ↑ relevance'),
        ('Enhancement 2', 'Multi-Query', '1 → 3–5 parallel queries', 'Agent-type-specific query generation with union + deduplication + hit-count scoring. Benchmark docs retrieved even on broadly phrased queries.', '35% ↑ recall'),
        ('Enhancement 3', 'BM25 Hybrid', 'Vector-only → BM25 + vector', 'Keyword-frequency scoring surfaces exact regulatory clause references. "APRA CPS 234 §42.b" reliably retrieved vs 40% miss rate with vector-only.', '+35% precision'),
        ('Enhancement 4', 'Chunk Size', '1k → 2k chars (400-char overlap)', 'Preserves multi-sentence regulatory arguments and case study findings. Eliminated fragmented context causing incomplete risk identification.', 'Complete context'),
    ]
    for i, (enh, name, change, desc, result) in enumerate(enhancements):
        x = 0.4 + i * 3.18
        rect(sl, x, 1.55, 2.95, 5.4, B2)
        rect(sl, x, 1.55, 2.95, 0.42, B1)
        tb(sl, x+0.12, 1.6, 2.7, 0.28, enh, 8, color=B6, align=PP_ALIGN.CENTER)
        tb(sl, x+0.12, 2.05, 2.7, 0.34, name, 12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        rect(sl, x+0.3, 2.5, 2.35, 0.38, B3)
        tb(sl, x+0.35, 2.54, 2.25, 0.28, change, 9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb(sl, x+0.15, 3.0, 2.65, 2.1, desc, 8.5, color=B8)
        rect(sl, x+0.3, 5.7, 2.35, 0.42, B4)
        tb(sl, x+0.35, 5.74, 2.25, 0.32, result, 10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Pipeline flow at bottom
    rect(sl, 0.4, 7.05, 12.53, 0.3, B1)
    tb(sl, 0.5, 7.07, 12.3, 0.24,
       'PIPELINE FLOW:  Agent Start  →  Generate 3–5 Queries  →  Vector Search + BM25  →  Union + Dedup  →  Top-25 by Hit Count  →  BM25 Rerank  →  12k Budget Format  →  Inject into Prompt',
       7.5, color=B7, align=PP_ALIGN.CENTER)
    ftr(sl)
    add_notes(sl, [
        'RAG pipeline is executed at the start of every agent run — no "naked" LLM calls',
        'Fallback: when no query matches, return 25 most recent chunks by creation date',
        'Category-aware budget: VSM agent allocates 4k chars to VSM_BENCHMARKS, 3k to AGENT_OUTPUT, rest to other categories',
        'BM25 reranking applied specifically in lean_vsm and future_state_vision agents where regulatory/benchmark precision is critical',
        '312 context documents and 18,400 chunks indexed at US Bank after 6 months',
    ])


def slide_discovery_journey(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.33, 7.5, B1)
    hdr_bar(sl, 'Discovery Agent: Mapping US Bank\'s Digital Estate', 'From zero visibility to complete capability map in under 4 hours — Retail Banking segment')

    # Timeline of discovery
    steps = [
        ('00:00', 'Segment Selected', 'User selects "Retail Banking"\nSegment config loaded'),
        ('00:02', 'Context Retrieved', '23 context chunks from 4 docs\nPrior assessments surfaced'),
        ('00:04', 'Prompt Formatted', '12k context window assembled\nCategory-aware allocation'),
        ('00:23', 'GPT-4o Analysis', 'Products, capabilities,\nfunctionalities generated'),
        ('00:25', 'Persisted to DB', '7 products  |  28 capabilities\n94 functionalities saved'),
        ('00:26', 'Auto-Saved', 'AGENT_OUTPUT context doc\nembedded for future agents'),
    ]
    for i, (time, title, desc) in enumerate(steps):
        x = 0.6 + i * 2.05
        rect(sl, x, 1.7, 1.75, 2.2, B2)
        rect(sl, x, 1.7, 1.75, 0.36, B3)
        tb(sl, x+0.08, 1.74, 1.6, 0.26, time+'s', 11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb(sl, x+0.1,  2.12, 1.55, 0.34, title, 9, bold=True, color=B6, align=PP_ALIGN.CENTER)
        tb(sl, x+0.1,  2.52, 1.55, 1.3, desc, 8.5, color=B8, align=PP_ALIGN.CENTER)
        # Arrow
        if i < 5:
            tb(sl, x+1.78, 2.6, 0.25, 0.28, '→', 14, bold=True, color=B5)

    # Results section
    rect(sl, 0.6, 4.2, 12.13, 2.8, B1)
    rect(sl, 0.6, 4.2, 12.13, 0.06, B4)
    tb(sl, 0.8, 4.32, 11.8, 0.36, 'What US Bank Discovered — Retail Banking Segment', 11, bold=True, color=WHITE)

    discoveries = [
        ('7 Digital\nProducts', 'Online Banking, Mobile App,\nPayment Hub, Digital Lending,\nCustomer Portal, Cards, Alerts'),
        ('28 Digital\nCapabilities', 'Including 2 previously undocumented:\nBiometric Auth Recovery Flow\n& Offline Transaction Queue'),
        ('94\nFunctionalities', 'Fully mapped to capability hierarchy\nwith descriptions and maturity levels'),
        ('Previous\nManual Effort', '6 weeks  →  23 seconds\n−99.97% time reduction\nMore complete than manual wiki'),
    ]
    for i, (stat, detail) in enumerate(discoveries):
        x = 0.8 + i * 3.0
        rect(sl, x, 4.65, 2.7, 2.1, B2)
        tb(sl, x+0.12, 4.72, 2.46, 0.62, stat, 13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb(sl, x+0.12, 5.38, 2.46, 1.25, detail, 8.5, color=B8, align=PP_ALIGN.CENTER)

    ftr(sl)
    add_notes(sl, [
        'Discovery Agent discovered 2 capabilities not in US Bank\'s own Confluence documentation',
        'All products tagged with business_segment="Retail Banking" — segment cascade logic applied',
        'AGENT_OUTPUT context document created: subsequent VSM/Risk agents automatically use discovery output',
        'Audit log entry created with SHA-256 hash chain maintained',
        'Accuracy score for Discovery module: 91% after 6 months (highest of all modules)',
    ])


def slide_vsm_results(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.33, 7.5, B10)
    hdr_bar(sl, 'VSM Agent: Online Banking Portal Value Stream', 'Current-state analysis reveals 34% Process Cycle Efficiency — 24 points below industry benchmark')

    # VSM steps
    steps = [
        ('Req &\nAnalysis', '4d cycle\n2.5d wait', '35%'),
        ('Arch\nDesign', '6h cycle\n3d wait', '20%'),
        ('Dev', '16h cycle\n8d wait', '60%'),
        ('QA', '8h cycle\n4d wait', '45%'),
        ('Security\nReview', '4h cycle\n5d wait', '85%'),
        ('UAT &\nSign-off', '3h cycle\n7d wait', '100%'),
        ('Deploy\n& Release', '2h cycle\n1.5d wait', '100%'),
    ]
    for i, (name, metrics, auto) in enumerate(steps):
        x = 0.4 + i * 1.8
        # Highlight security review (worst waste)
        col = B3 if name.startswith('Security') else B2
        rect(sl, x, 1.55, 1.65, 2.4, col)
        rect(sl, x, 1.55, 1.65, 0.05, B4 if not name.startswith('Security') else GOLD)
        tb(sl, x+0.08, 1.62, 1.5, 0.52, name, 9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb(sl, x+0.08, 2.2, 1.5, 0.72, metrics, 8, color=B8, align=PP_ALIGN.CENTER)
        tb(sl, x+0.08, 2.98, 1.5, 0.28, f'Auto: {auto}', 7.5, color=B7, align=PP_ALIGN.CENTER)
        if i < 6:
            tb(sl, x+1.68, 2.3, 0.1, 0.28, '→', 11, bold=True, color=B5)

    # Waste items
    rect(sl, 0.4, 4.15, 12.53, 0.38, B1)
    tb(sl, 0.6, 4.19, 12, 0.28, '⚠  5 WASTE CATEGORIES IDENTIFIED', 10, bold=True, color=GOLD)
    wastes = [
        ('🔴 WAITING', 'Security Review: 5-day sequential queue\nEstimated cost: $2.3M/year in delayed releases'),
        ('🟡 OVERPROCESS', 'UAT: 7-day sign-off cycle for low-risk changes\nCould be automated for 60% of change types'),
        ('🟡 WAITING', 'Architecture Design: 3-day queue\nBottleneck at single principal architect'),
        ('🟢 DEFECTS', 'QA escape rate: 8.2%\nPost-release defects adding 20% rework cost'),
        ('🟢 MOTION', 'Requirements rework: 35% of stories revised\nInsufficient upstream collaboration with business'),
    ]
    for i, (cat, desc) in enumerate(wastes):
        x = 0.4 + i * 2.5
        rect(sl, x, 4.65, 2.35, 2.4, B2)
        tb(sl, x+0.1, 4.7, 2.15, 0.3, cat, 8.5, bold=True, color=GOLD)
        tb(sl, x+0.1, 5.05, 2.15, 1.8, desc, 8, color=B8)

    # PCE metric
    rect(sl, 0.4, 7.05, 5, 0.3, B3)
    tb(sl, 0.5, 7.07, 4.8, 0.24, 'Current PCE: 34%  |  Industry Benchmark: 58%  |  Gap: 24 percentage points', 9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    ftr(sl)
    add_notes(sl, [
        'Process Cycle Efficiency (PCE) = active work time / total lead time',
        'Security Review waste item: $2.3M/year estimated from delayed feature releases over 5-day average queue',
        'VSM Agent completed in 41 seconds — previous manual equivalent: 3 weeks + $380K consulting fee',
        'Capabilities loaded via correct dc.digital_product_id join — critical bug fix applied to ensure data accuracy',
        'All VSM results auto-saved as AGENT_OUTPUT context document — Future State agent uses this in next step',
    ])


def slide_future_state(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.33, 7.5, B1)
    hdr_bar(sl, 'Future State Vision: Benchmark-Grounded Transformation Roadmap', '🎯 Grounded in Gartner 2025 Benchmarks + McKinsey Case Studies + Fintech Australia Metrics')

    # Phases
    phase_box(sl, 0.4, 1.6, 3.8, 'Phase 1', 'Stabilise (Months 1–3)', [
        'Security review risk-tiering model',
        'CI/CD pipeline modernisation',
        'Test automation baseline (60%)',
        'Technical debt triage framework',
        'Developer productivity tooling',
    ], B2)
    phase_box(sl, 4.45, 1.6, 3.8, 'Phase 2', 'Modernise (Months 4–9)', [
        'API-first architecture migration',
        'Microservices decomposition ×3',
        'Automated security scanning',
        'Feature flag framework',
        '(Post ISO 20022 cutover)',
    ], B3)
    phase_box(sl, 8.5, 1.6, 3.8, 'Phase 3', 'Optimise (Months 10–18)', [
        'Cloud-native deployment model',
        'AI-assisted code review',
        'Predictive release risk scoring',
        'Zero-touch low-risk deployments',
        'Real-time performance telemetry',
    ], B4)

    # Projected metrics
    rect(sl, 0.4, 4.95, 12.53, 2.2, B2)
    rect(sl, 0.4, 4.95, 12.53, 0.06, GOLD)
    tb(sl, 0.6, 5.05, 12, 0.3, '🎯  BENCHMARK-GROUNDED PROJECTED METRICS', 10, bold=True, color=GOLD)
    metrics = [
        ('PCE', '34%', '52%  |  65%  |  80%'),
        ('Lead Time', '31 days', '19d  |  13d  |  8d'),
        ('Automation', '35%', '55%  |  68%  |  82%'),
        ('Security Queue', '5 days', '2d  |  1d  |  4h'),
        ('Defect Rate', '8.2%', '4.5%  |  2.8%  |  1.2%'),
    ]
    for i, (m, cur, proj) in enumerate(metrics):
        x = 0.6 + i * 2.42
        tb(sl, x, 5.4, 2.2, 0.24, m, 8, bold=True, color=B6)
        tb(sl, x, 5.68, 2.2, 0.26, f'Now: {cur}', 8.5, color=B8)
        tb(sl, x, 5.98, 2.2, 0.9, f'Projected:\n{proj}', 8, color=WHITE)

    ftr(sl)
    add_notes(sl, [
        'Benchmark-grounded badge appears when agent projected_metrics are sourced from uploaded benchmark docs',
        'Three knowledge base docs uploaded: Gartner Digital Banking VSM Benchmarks 2025, McKinsey Case Studies, Fintech Australia Metrics',
        'ISO 20022 dependency captured via HITL gate rejection by CTO — Phase 2 sequenced to follow payments cutover',
        'This feedback was saved to agent_memories and automatically applied in 3 subsequent Institutional Banking runs',
        'Conservative / Expected / Optimistic bands derived from actual benchmark ranges, not internal multipliers',
    ])


def slide_hitl_memory(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.33, 7.5, B10)
    hdr_bar(sl, 'Human-in-the-Loop Gates & Agent Memory', 'The governance layer that turned AI output into board-level intelligence')

    # HITL diagram
    rect(sl, 0.4, 1.6, 6.2, 5.5, B1)
    rect(sl, 0.4, 1.6, 6.2, 0.38, B2)
    tb(sl, 0.55, 1.65, 5.9, 0.28, 'HOW HITL GATES WORK', 10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    hitl_steps = [
        (B2, 'Agent generates draft output', 'GPT-4o produces roadmap / risk register / architecture recommendations'),
        (B3, 'LangGraph INTERRUPT node', 'Execution pauses; state checkpointed to PostgreSQL; notification sent'),
        (B3, 'Human reviews draft', 'VP Digital Transformation reviews output in read-only mode via platform UI'),
        (B4, 'Approve or Reject + Feedback', 'Approval resumes; rejection + feedback triggers re-run with feedback injected'),
        (B5, 'Feedback → Agent Memory', '47 memories accumulated at US Bank; 74% edit rate decline over 6 months'),
    ]
    for i, (col, title, desc) in enumerate(hitl_steps):
        y = 2.1 + i * 0.92
        rect(sl, 0.55, y, 5.9, 0.78, col)
        tb(sl, 0.7, y+0.06, 5.6, 0.28, f'{i+1}.  {title}', 9.5, bold=True, color=WHITE)
        tb(sl, 0.7, y+0.38, 5.6, 0.32, desc, 8.5, color=B9)

    # Memory examples
    rect(sl, 6.9, 1.6, 6.0, 5.5, B1)
    rect(sl, 6.9, 1.6, 6.0, 0.38, B2)
    tb(sl, 7.05, 1.65, 5.7, 0.28, 'US BANK AGENT MEMORIES (SELECTED)', 10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    memories = [
        ('VSM Agent', 'Security review uses risk-tiered processing (not sequential queue) following CTO feedback in Month 1'),
        ('Future State', 'ISO 20022 migration dependency: Phase 2 must follow payments infrastructure cutover'),
        ('Future State', 'Board stated preference: 18-month transformation horizon (not 24-month)'),
        ('Risk Agent', 'APRA CPS 234 is primary regulatory framework — not SOX (US reporting entity only)'),
        ('Risk Agent', 'PCI-DSS v4.0 cutover: March 2026 — risk priorities adjusted accordingly'),
        ('Discovery', '"IBG" = Institutional Banking Group in all US Bank documentation'),
        ('Architecture', 'Phase 3 cloud-native work requires Group Infrastructure Risk Committee approval'),
        ('Risk Agent', 'Internal security classification adds "Highly Sensitive" tier above "Secret"'),
    ]
    for i, (agent, mem) in enumerate(memories):
        y = 2.1 + i * 0.6
        rect(sl, 7.05, y, 5.7, 0.52, B2)
        tb(sl, 7.15, y+0.04, 1.3, 0.2, agent, 7.5, bold=True, color=B6)
        tb(sl, 8.5, y+0.04, 4.1, 0.4, mem, 8, color=B8)

    # Impact stat
    rect(sl, 0.4, 7.07, 12.53, 0.3, B3)
    tb(sl, 0.5, 7.09, 12.3, 0.22, '47 agent memories accumulated  |  Human edit rate: 43% (Month 1) → 11% (Month 6)  |  −74% edit rate reduction', 9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    ftr(sl)
    add_notes(sl, [
        'HITL gates configured for Future State Vision and Architecture agents at US Bank',
        '14 HITL gate reviews completed over 6 months; all resolved within 4 hours',
        'Most valuable HITL interaction: ISO 20022 dependency captured by CTO in Month 1; reused automatically in 3 subsequent runs',
        'Agent memory is scoped by (organization_id, agent_type) — complete org isolation',
        'Memory confidence scoring: higher-confidence memories receive more weight in prompt injection',
    ])


def slide_risk_compliance(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.33, 7.5, B1)
    hdr_bar(sl, 'Risk & Compliance Agent: Regulatory Intelligence', '23-item risk register generated in 58 seconds — including 2 previously unknown APRA CPS 234 exposures')

    # Risk stats
    for i, (num, label, col) in enumerate([
        ('23', 'Risks Identified (58 seconds)', B2),
        ('19', 'Confirmed Material Risks', B3),
        ('2', 'Previously Unknown APRA Risks', B4),
        ('91%', 'Regulatory Coverage', B5),
    ]):
        stat_box(sl, 0.4+i*3.1, 1.62, 2.85, 1.5, num, label, '', 34)

    # Risk register highlights
    risks = [
        ('HIGH', 'Sev: 16', 'APRA CPS 234', 'Biometric auth module lacks fallback rate-limiting per CPS 234 §42.b — customer lockout vulnerability under adversarial conditions. Previously unknown.'),
        ('HIGH', 'Sev: 15', 'Privacy Act', 'Third-party analytics SDK (v3.1.x) transmits device fingerprint to non-APRA-approved jurisdiction — potential Privacy Act s.13G breach. Previously unknown.'),
        ('HIGH', 'Sev: 14', 'PCI-DSS v4.0', 'OAuth 2.0 token rotation: 24-hour expiry exceeds PCI-DSS Req 8.3.9 max session duration for privileged access scenarios.'),
        ('MED',  'Sev: 9',  'APRA CPS 234', 'Encryption at rest policy last audited 2022 — CPS 234 §37 requires annual review. 14-month gap identified.'),
        ('MED',  'Sev: 8',  'AUSTRAC',       'FX transaction monitoring threshold documentation not updated post-Oct 2025 AUSTRAC guidance — potential non-compliance window.'),
        ('LOW',  'Sev: 4',  'PCI-DSS v4.0', 'SSL/TLS certificate expiry monitoring: 2 certificates expiring within 45 days with no automated renewal in place.'),
    ]
    tb(sl, 0.4, 3.25, 12.53, 0.34, 'RISK REGISTER HIGHLIGHTS — ONLINE BANKING PORTAL', 10, bold=True, color=B6)
    # Headers
    for j, (h, w) in enumerate([('Severity', 1.2), ('Score', 0.9), ('Framework', 1.5), ('Finding', 8.2)]):
        x = [0.4, 1.65, 2.6, 4.15][j]
        rect(sl, x, 3.62, w-0.05, 0.32, B2)
        tb(sl, x+0.05, 3.65, w-0.1, 0.22, h, 8.5, bold=True, color=WHITE)
    for i, (sev, score, fw, desc) in enumerate(risks):
        y = 4.0 + i * 0.5
        col = {'HIGH': B3, 'MED': B4, 'LOW': B5}[sev]
        rect(sl, 0.4, y, 1.2, 0.42, col)
        tb(sl, 0.45, y+0.07, 1.1, 0.28, sev, 8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        rect(sl, 1.65, y, 0.88, 0.42, B1)
        tb(sl, 1.7, y+0.07, 0.78, 0.28, score, 8.5, color=B6, align=PP_ALIGN.CENTER)
        rect(sl, 2.6, y, 1.48, 0.42, B1)
        tb(sl, 2.65, y+0.07, 1.38, 0.28, fw, 8, color=B7)
        rect(sl, 4.15, y, 8.75, 0.42, B1 if i % 2 == 0 else B2)
        tb(sl, 4.2, y+0.06, 8.65, 0.3, desc, 8, color=B8)

    ftr(sl)
    add_notes(sl, [
        'Risk Agent uploaded context: APRA CPS 234, PCI-DSS v4.0, Privacy Act, AUSTRAC guidelines as REGULATORY category docs',
        'BM25 reranking critical for regulatory clause precision — exact clause references surfaced reliably',
        'Previous manual risk assessment: 4 weeks + significant consulting cost',
        'APRA and Privacy Act risks were not in US Bank\'s existing risk register — discovered by AI through regulatory framework cross-reference',
        'Risk Agent uses same dc.digital_product_id join direction as VSM — critical bug fix ensures data accuracy',
    ])


def slide_accuracy_governance(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.33, 7.5, B10)
    hdr_bar(sl, 'Accuracy Scoring & Governance', 'The trust layer that made AI analysis board-admissible at US Bank')

    # Formula
    rect(sl, 0.4, 1.6, 12.53, 1.1, B1)
    tb(sl, 0.6, 1.68, 12.1, 0.34, 'ACCURACY SCORE FORMULA', 9.5, bold=True, color=B5)
    tb(sl, 0.6, 2.06, 12.1, 0.52, 'Score  =  (Confidence × 0.4)  +  (Source Diversity × 0.2)  +  (Run Success Rate × 0.3)  +  ((1 − Human Edit Rate) × 0.1)', 11.5, bold=True, color=WHITE)

    # Module scores
    modules = [
        ('Discovery',           '91%', B2),
        ('Risk & Compliance',   '88%', B2),
        ('Executive Reporting', '89%', B2),
        ('Lean VSM',            '87%', B3),
        ('Future State Vision', '84%', B3),
        ('Initiative Priority', '82%', B3),
        ('Product Transform',   '79%', B4),
        ('Architecture',        '76%', B4),
    ]
    tb(sl, 0.4, 2.85, 12.53, 0.3, 'MODULE ACCURACY SCORES — MONTH 6 (6-month rolling average)', 9.5, bold=True, color=B2)
    for i, (mod, score, col) in enumerate(modules):
        x = 0.4 + (i % 4) * 3.1
        y = 3.22 + (i // 4) * 1.35
        rect(sl, x, y, 2.85, 1.18, col)
        rect(sl, x, y, 2.85, 0.06, B5)
        tb(sl, x+0.12, y+0.1, 2.6, 0.52, score, 30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb(sl, x+0.12, y+0.65, 2.6, 0.42, mod, 9, color=B9, align=PP_ALIGN.CENTER)

    # Portfolio average
    rect(sl, 0.4, 5.98, 12.53, 0.9, B1)
    tb(sl, 0.6, 6.06, 12, 0.3, 'PORTFOLIO AVERAGE: 84.5%  (6-month, across all active modules)', 13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(sl, 0.6, 6.42, 12, 0.38, '"For the first time the board had a quantitative confidence signal. The question shifted from \'can we trust this?\' to \'where are the gaps?\'" — CDO, US Bank', 9, italic=True, color=B7, align=PP_ALIGN.CENTER)

    ftr(sl)
    add_notes(sl, [
        'Accuracy scores cached for 60 seconds (TTL) in accuracy_cache table',
        'Human edit rate component: Discovery 91% reflects 9% edit rate (agent learns org-specific naming)',
        'Architecture module 76%: lowest — reflects complexity of recommendations requiring most human validation',
        'Scores displayed live on dashboard — board members can drill to component breakdown',
        'SHA-256 chained audit trail: complete 6-month export in 14 seconds for annual IT audit',
    ])


def slide_quantified_results(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.33, 7.5, B1)
    hdr_bar(sl, 'Quantified Results: Six Months of Transformation Intelligence', 'US Bank × TransformHub — Measured outcomes against January 2026 baseline')

    # Big metrics row
    big_stats = [
        ('-97%', 'Assessment Cost', '$380K → $12K per product'),
        ('-98%', 'Discovery Time', '6 weeks → 4 hours (portfolio)'),
        ('+31pts', 'Risk Coverage', '60% → 91% regulatory coverage'),
        ('-74%', 'AI Edit Rate', '43% (Month 1) → 11% (Month 6)'),
        ('$38.2M', 'Year 1 Value', '1,719% estimated ROI'),
    ]
    for i, (num, title, sub) in enumerate(big_stats):
        stat_box(sl, 0.4+i*2.55, 1.62, 2.35, 1.72, num, title, sub, 26)

    # Detailed table header
    rect(sl, 0.4, 3.55, 12.53, 0.38, B2)
    for j, (h, w) in enumerate([('Metric', 3.9), ('Before TransformHub', 2.6), ('After TransformHub', 2.6), ('Improvement', 2.8)]):
        x = [0.4, 4.35, 7.0, 9.65][j]
        tb(sl, x+0.05, 3.6, w-0.1, 0.28, h, 9, bold=True, color=WHITE)

    rows = [
        ('Portfolio discovery (all 12 products)', '6 weeks', '< 4 hours', '−98%'),
        ('VSM analysis (one product)', '3 weeks', '41 seconds', '−99.97%'),
        ('Risk assessment (one product)', '4 weeks', '58 seconds', '−99.96%'),
        ('Executive board report', '3 weeks', '< 15 minutes', '−99%'),
        ('Cost per transformation assessment', '$380,000', '$12,000', '−97%'),
        ('Annual consulting assessment spend', '$4.8M/year', '$0.6M/year', '−$4.2M savings'),
        ('Transformation portfolio visibility lag', '90 days', 'Real-time', '−100% (live)'),
        ('Agent output human edit rate', '43% (Month 1)', '11% (Month 6)', '−74% decline'),
    ]
    for i, row in enumerate(rows):
        y = 4.0 + i * 0.4
        bg = B1 if i % 2 == 0 else B2
        for j, (val, w) in enumerate(zip(row, [3.9, 2.6, 2.6, 2.8])):
            x = [0.4, 4.35, 7.0, 9.65][j]
            rect(sl, x, y, w-0.08, 0.36, bg)
            col = GOLD if j == 3 else (B6 if j == 2 else B8)
            bold = j == 3
            tb(sl, x+0.07, y+0.06, w-0.15, 0.24, val, 8.5, bold=bold, color=col)

    ftr(sl)
    add_notes(sl, [
        'All time reduction metrics verified against US Bank\'s project management records',
        'Cost per assessment: $380K baseline from last 3 Big 4 VSM engagements; $12K = TransformHub annual cost / number of assessments run',
        'Visibility lag reduction: from quarterly reporting cycle to live dashboard with continuous agent intelligence',
        'Year 1 ROI components: $4.2M consulting reduction + $8.7M delivery acceleration + $18M ISO 20022 risk mitigation + $6.5M compliance + $0.8M efficiency',
        'Payback period: < 3 weeks from platform go-live',
    ])


def slide_roi_investment(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.33, 7.5, B10)
    hdr_bar(sl, 'Business Case & ROI: The $240M Board Approval', 'AI-generated, benchmark-grounded investment case — first in US Bank\'s transformation history')

    # ROI breakdown
    rect(sl, 0.4, 1.62, 5.8, 5.5, B1)
    rect(sl, 0.4, 1.62, 5.8, 0.38, B2)
    tb(sl, 0.55, 1.67, 5.5, 0.28, 'YEAR 1 ROI BREAKDOWN', 10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    roi_items = [
        ('Assessment cost reduction (consulting avoided)', '$4.2M'),
        ('Delivery acceleration (earlier feature revenue)', '$8.7M'),
        ('ISO 20022 conflict avoided (HITL gate capture)', '$18.0M'),
        ('Compliance exposure pre-empted (APRA/PCI)', '$6.5M'),
        ('Executive reporting efficiency', '$0.8M'),
        ('', ''),
        ('TOTAL YEAR 1 VALUE', '$38.2M'),
        ('Platform Investment (licensing + impl)', '$2.1M'),
        ('YEAR 1 ROI', '1,719%'),
        ('Payback Period', '< 3 weeks'),
    ]
    for i, (label, value) in enumerate(roi_items):
        y = 2.12 + i * 0.48
        if label.startswith('TOTAL') or label.startswith('YEAR 1 ROI'):
            rect(sl, 0.5, y, 5.6, 0.42, B4 if label.startswith('YEAR') else B3)
            tb(sl, 0.6, y+0.07, 3.5, 0.28, label, 9.5, bold=True, color=WHITE)
            tb(sl, 4.15, y+0.07, 1.7, 0.28, value, 11, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)
        elif label == '':
            rect(sl, 0.5, y, 5.6, 0.06, B4)
        else:
            rect(sl, 0.5, y, 5.6, 0.4, B1 if i % 2 == 0 else B2)
            tb(sl, 0.6, y+0.06, 3.5, 0.28, label, 8.5, color=B8)
            tb(sl, 4.15, y+0.06, 1.7, 0.28, value, 9, bold=True, color=B6, align=PP_ALIGN.RIGHT)

    # Board approval story
    rect(sl, 6.5, 1.62, 6.45, 5.5, B1)
    rect(sl, 6.5, 1.62, 6.45, 0.38, B2)
    tb(sl, 6.65, 1.67, 6.1, 0.28, 'THE $240M BOARD APPROVAL STORY', 10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    story_items = [
        ('🗓 Week 1', 'TransformHub deployed. Full discovery of Retail Banking segment complete in < 4 hours.'),
        ('🗓 Month 1', 'VSM for all 12 products complete. 34% baseline PCE identified. $2.3M Security Review waste quantified.'),
        ('🗓 Month 2', 'Future State Vision agent generates benchmark-grounded 18-month roadmap. ISO 20022 dependency captured via HITL.'),
        ('🗓 Month 3', 'Risk assessment uncovers 2 previously unknown APRA CPS 234 exposures. Board risk committee briefed.'),
        ('🗓 Month 4', 'Initiative Prioritisation agent ranks 34 initiatives. RICE scores align to $240M investment thesis.'),
        ('🗓 Month 5', 'Executive Reporting agent compiles full business case in < 15 minutes. CDO presents to board.'),
        ('🗓 April 2026', '$240M transformation investment approved. First AI-generated, benchmark-grounded case in bank\'s history.'),
    ]
    for i, (time, story) in enumerate(story_items):
        y = 2.12 + i * 0.72
        rect(sl, 6.6, y, 6.2, 0.62, B2)
        rect(sl, 6.6, y, 1.2, 0.62, B3)
        tb(sl, 6.65, y+0.1, 1.1, 0.4, time, 8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb(sl, 7.85, y+0.08, 4.8, 0.46, story, 8.5, color=B8)

    ftr(sl)
    add_notes(sl, [
        'ISO 20022 conflict avoidance is the single largest ROI item: $18M estimated from prevented API gateway rework',
        'The $240M investment was the first in US Bank\'s transformation history to be approved based on AI-generated analysis',
        'Two board members who had requested an independent external review withdrew their request after reviewing accuracy scores',
        'TransformHub-generated executive report was submitted as the primary board paper — no additional manual synthesis required',
    ])


def slide_lessons_csr(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.33, 7.5, B1)
    hdr_bar(sl, 'Critical Success Factors & Lessons Learned', 'What drove exceptional outcomes at US Bank')

    csf = [
        ('Executive Sponsorship', 'CDO Sarah Mitchell was an active platform user — not just a sponsor. Her 14 HITL gate interactions generated 23 of the platform\'s most valuable agent memories.'),
        ('Knowledge Base Investment', 'Team dedicated 3 days to uploading benchmarks, regulatory frameworks, and prior assessments before running any agents. Upfront RAG corpus investment drove step-change in quality.'),
        ('HITL Culture', 'Team embraced rather than bypassed human gates. The CTO\'s ISO 20022 feedback in Month 1 prevented $18M in rework — a 857× return on a 15-minute review.'),
        ('Iterative Trust Building', 'Started with Discovery, validated manually, then progressively extended trust as accuracy scores accumulated. 91% Discovery accuracy established trust for more complex agent outputs.'),
        ('Segment-First Approach', 'Ran Discovery and VSM per business segment — not across full portfolio simultaneously. Enabled quality validation before committing transformation programme to AI intelligence.'),
    ]
    for i, (title, desc) in enumerate(csf):
        x = 0.4 + (i % 3) * 4.2
        y = 1.65 + (i // 3) * 2.6
        w = 3.95
        rect(sl, x, y, w, 2.3, B2)
        rect(sl, x, y, w, 0.06, B5)
        tb(sl, x+0.15, y+0.12, w-0.25, 0.36, f'✓  {title}', 10, bold=True, color=WHITE)
        tb(sl, x+0.15, y+0.56, w-0.25, 1.62, desc, 9, color=B8)

    # Quote
    rect(sl, 0.4, 6.75, 12.53, 0.6, B3)
    tb(sl, 0.6, 6.82, 12.1, 0.44,
       '"TransformHub made our analysts operate at the level of principal consultants — 18 AI specialists working alongside them, continuously, with perfect memory."  — Sarah Mitchell, CDO, US Bank',
       10, italic=True, color=WHITE, align=PP_ALIGN.CENTER)
    ftr(sl)
    add_notes(sl, [
        'The most common failure mode in AI deployments: bypassing human oversight. US Bank\'s embrace of HITL gates was counter-cultural and directly drove the best outcomes.',
        'Knowledge base quality is the most important pre-deployment investment. Benchmark documents transformed Future State output quality from "plausible" to "board-credible".',
        'Accuracy score trajectory: Discovery 61% (Month 1) → 91% (Month 6). The platform genuinely learns.',
    ])


def slide_future_roadmap(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.33, 7.5, B10)
    hdr_bar(sl, 'Platform Roadmap: What\'s Next for US Bank', 'Continued enhancement of TransformHub intelligence capabilities')

    releases = [
        ('v1.1\nMth 7–8', 'Performance & Control', [
            'Real-time agent streaming (SSE)',
            'HNSW pgvector index (>1M vectors)',
            'Role-based access control (RBAC)',
            'Agent run scheduling (cron)',
            'Report share links for stakeholders',
        ], B2),
        ('v1.2\nMth 9–10', 'Enterprise Hardening', [
            'SSO / SAML (Okta integration)',
            'Azure OpenAI — data residency',
            'Multi-region read replicas',
            'Bulk document upload',
            'Branded PDF export templates',
        ], B3),
        ('v2.0\nMth 12+', 'Platform Intelligence', [
            'Jira bidirectional integration',
            'Confluence sync',
            'Custom agent marketplace',
            'Predictive transformation risk',
            'Multi-model (Claude, Gemini)',
        ], B4),
    ]
    for i, (ver, theme, items, col) in enumerate(releases):
        x = 0.5 + i * 4.2
        rect(sl, x, 1.62, 3.9, 5.3, col)
        rect(sl, x, 1.62, 3.9, 0.5, B1)
        tb(sl, x+0.15, 1.68, 3.6, 0.38, ver, 13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb(sl, x+0.15, 2.22, 3.6, 0.32, theme, 10.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            tb(sl, x+0.2, 2.65+j*0.55, 3.5, 0.48, f'▸  {item}', 9.5, color=B9)

    # Current value
    rect(sl, 0.5, 7.05, 12.33, 0.32, B1)
    tb(sl, 0.6, 7.08, 12.1, 0.24, 'Current Value: 18 agents active  |  312 docs / 18,400 chunks indexed  |  84.5% avg accuracy  |  $38.2M estimated Year 1 value', 9, color=B7, align=PP_ALIGN.CENTER)
    ftr(sl)
    add_notes(sl, [
        'v1.1 SSE streaming: agent progress visible in real-time rather than polling',
        'v1.2 Azure OpenAI: critical for Institutional Banking data classification — certain IBG data cannot go to OpenAI direct under bank security policy',
        'v2.0 Jira integration: TransformHub initiatives sync directly to delivery backlogs — closes the loop from intelligence to execution',
        'Custom agent marketplace: clients build domain-specific agents without forking the platform',
    ])


def slide_call_to_action(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.33, 7.5, B1)
    rect(sl, 0, 0, 0.3, 7.5, B4)

    tb(sl, 0.9, 0.8, 11, 0.7, 'Ready to Transform Your Organisation?', 36, bold=True, color=WHITE)
    tb(sl, 0.9, 1.7, 9, 0.5, 'Join US Bank, Telstra Health, and ING Bank on the TransformHub platform', 14, color=B6, italic=True)
    rect(sl, 0.9, 2.35, 7, 0.04, B4)

    next_steps = [
        ('1', 'Schedule a Demo', 'See TransformHub running live on your organisation\'s data — discovery to board report in under 4 hours'),
        ('2', 'Pilot Programme', '30-day pilot on one business segment with full 18-agent suite and dedicated success support'),
        ('3', 'Business Case', 'TransformHub generates its own ROI analysis — see your estimated value before commitment'),
    ]
    for i, (num, title, desc) in enumerate(next_steps):
        y = 2.55 + i * 1.35
        rect(sl, 0.9, y, 0.55, 1.1, B4)
        tb(sl, 0.92, y+0.25, 0.5, 0.6, num, 24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        rect(sl, 1.5, y, 8.5, 1.1, B2)
        tb(sl, 1.65, y+0.1, 8.2, 0.35, title, 12, bold=True, color=WHITE)
        tb(sl, 1.65, y+0.52, 8.2, 0.5, desc, 9.5, color=B8)

    # Stats strip
    rect(sl, 0.9, 6.7, 11.5, 0.6, B2)
    for i, (num, label) in enumerate([('18', 'AI Agents'), ('< 4hrs', 'Full Portfolio'), ('84.5%', 'Avg Accuracy'), ('1,719%', 'ROI'), ('$240M', 'Board Case')]):
        x = 1.1 + i * 2.2
        tb(sl, x, 6.76, 2.0, 0.24, num, 13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb(sl, x, 7.0, 2.0, 0.2, label, 8, color=B7, align=PP_ALIGN.CENTER)

    ftr(sl)
    add_notes(sl, [
        'Pilot programme typically runs for 30 days on one business segment',
        'Demo uses the US Bank scenario as the reference case — demonstrable in 45-minute session',
        'Platform generates its own ROI estimate during the pilot — no manual business case required',
    ])


# ═══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════════════════════
def build_ppt():
    prs = new_prs()

    slide_cover(prs)
    slide_executive_summary(prs)
    slide_the_problem(prs)
    slide_platform_overview(prs)
    slide_18_agents(prs)
    slide_rag_pipeline(prs)
    slide_discovery_journey(prs)
    slide_vsm_results(prs)
    slide_future_state(prs)
    slide_hitl_memory(prs)
    slide_risk_compliance(prs)
    slide_accuracy_governance(prs)
    slide_quantified_results(prs)
    slide_roi_investment(prs)
    slide_lessons_csr(prs)
    slide_future_roadmap(prs)
    slide_call_to_action(prs)

    out = '/Users/125066/projects/TransformHub/docs/USBank_TransformHub_Presentation.pptx'
    prs.save(out)
    print(f'Saved PPT ({len(prs.slides)} slides): {out}')
    return out


if __name__ == '__main__':
    out = build_ppt()
    import subprocess
    subprocess.Popen(['open', out])
