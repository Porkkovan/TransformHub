"""
TransformHub — Consulting Service Offering Deck
26 slides · Story-driven flow · White background · Blue monochrome palette
Talking points: PPT Notes pane only (not on slides)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Blue monochrome palette (darkest → lightest) ───────────────────────────
#   B1  #0C1C3A  — darkest navy (headers, title bg)
#   B2  #1E3A8A  — dark blue    (phase 1 / critical emphasis)
#   B3  #1D4ED8  — medium-dark  (phase 2 / primary accent)
#   B4  #2563EB  — medium blue  (phase 3 / main interactive)
#   B5  #3B82F6  — medium-light (phase 4 / secondary)
#   B6  #60A5FA  — light blue   (phase 5 / soft accent)
#   B7  #93C5FD  — lighter      (decorative / borders)
#   B8  #BFDBFE  — very light   (subtle highlights)
#   B9  #DBEAFE  — bg tint dark (card backgrounds)
#   B10 #EFF6FF  — bg tint light (alternate row)
#   B11 #F8FBFF  — near-white blue (lightest bg)

B1  = RGBColor(0x0C, 0x1C, 0x3A)   # darkest navy
B2  = RGBColor(0x1E, 0x3A, 0x8A)   # dark blue
B3  = RGBColor(0x1D, 0x4E, 0xD8)   # medium-dark blue
B4  = RGBColor(0x25, 0x63, 0xEB)   # medium blue
B5  = RGBColor(0x3B, 0x82, 0xF6)   # medium-light blue
B6  = RGBColor(0x60, 0xA5, 0xFA)   # light blue
B7  = RGBColor(0x93, 0xC5, 0xFD)   # lighter blue
B8  = RGBColor(0xBF, 0xDB, 0xFE)   # very light blue
B9  = RGBColor(0xDB, 0xEA, 0xFE)   # bg tint (darker)
B10 = RGBColor(0xEF, 0xF6, 0xFF)   # bg tint (lighter)
B11 = RGBColor(0xF0, 0xF7, 0xFF)   # near-white blue

# Semantic aliases — 6 distinct steps across the blue scale
# (used as phase / category colours throughout all slides)
NAVY    = B1    # darkest — headers, title backgrounds
RED     = B2    # dark blue     — Phase 1 / critical / "before" column
AMBER   = B3    # medium-dark   — Phase 2 / secondary emphasis
BLUE    = B4    # medium        — Phase 3 / primary interactive
PURPLE  = B5    # medium-light  — Phase 4 / soft accent
GREEN   = B6    # light blue    — Phase 5 / positive outcomes
CYAN    = B7    # lighter blue  — Phase 6 / decorative

DARK    = RGBColor(0x0F, 0x17, 0x2A)
MID     = RGBColor(0x47, 0x55, 0x69)
LGRAY   = RGBColor(0xF1, 0xF5, 0xF9)
MGRAY   = RGBColor(0xE2, 0xE8, 0xF0)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

# Background tint aliases
REDBG   = B10   # was red bg   → lightest blue tint
AMBERBG = B10   # was amber bg → lightest blue tint
GREENBG = B11   # was green bg → near-white blue
BLUEBG  = B9    # was blue bg  → mid blue tint
NAVYBG  = B11

# ── Layout zones ───────────────────────────────────────────────────────────
CT   = 1.15   # content top
CB   = 6.88   # content bottom (full height, no TP strip)
FT   = 7.18   # footer top

W = Inches(13.33)
H = Inches(7.50)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── Core helpers ───────────────────────────────────────────────────────────
def ns():
    sl = prs.slides.add_slide(BLANK)
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = WHITE
    return sl

def rect(sl, x, y, w, h, fill=None, line=None, lw=Pt(0.75)):
    sh = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid() if fill else sh.fill.background()
    if fill: sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = lw
    else:    sh.line.fill.background()
    return sh

def oval(sl, x, y, w, h, fill):
    sh = sl.shapes.add_shape(9, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.fill.background()

def tb(sl, text, x, y, w, h, sz=10, bold=False, color=DARK,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return box

def hdr(sl, title, sub=None, accent=NAVY):
    rect(sl, 0, 0, 13.33, 0.10, fill=accent)
    tb(sl, title, 0.4, 0.14, 11.5, 0.60, sz=22, bold=True, color=DARK)
    if sub:
        tb(sl, sub, 0.4, 0.75, 12.5, 0.30, sz=10, color=MID, italic=True)

def ftr(sl, txt="TransformHub · AI-Powered Digital Transformation Intelligence  ·  Consulting Service Offering"):
    rect(sl, 0, FT, 13.33, 0.30, fill=LGRAY)
    tb(sl, txt, 0.35, FT + 0.05, 12.6, 0.20, sz=10, color=MID)

def add_notes(sl, pts):
    nf = sl.notes_slide.notes_text_frame; nf.clear()
    p = nf.paragraphs[0]; r = p.add_run()
    r.text = "TALKING POINTS"; r.font.bold = True; r.font.size = Pt(13)
    for pt in pts:
        para = nf.add_paragraph(); run = para.add_run()
        run.text = f"•  {pt}"; run.font.size = Pt(11)

def tp(sl, pts):
    """Write talking points to Notes pane only — no visual strip on slide."""
    add_notes(sl, pts)

def metric_box(sl, val, lbl, x, y, w=2.8, h=1.0, vc=NAVY, bg=LGRAY, bc=NAVY):
    rect(sl, x, y, w, h, fill=bg, line=bc, lw=Pt(0.8))
    tb(sl, val, x, y + 0.06, w, 0.52, sz=28, bold=True, color=vc, align=PP_ALIGN.CENTER)
    tb(sl, lbl, x, y + 0.60, w, 0.34, sz=10, color=MID, align=PP_ALIGN.CENTER)

def row_label(sl, lbl, val, x, y, w=12.3, lc=DARK, vc=BLUE):
    tb(sl, lbl, x, y, w * 0.62, 0.24, sz=10, color=lc)
    tb(sl, val, x + w * 0.63, y, w * 0.36, 0.24, sz=10, bold=True, color=vc, align=PP_ALIGN.RIGHT)

def phase_card(sl, num, title, color, items, x, y, w=2.4, h=3.6):
    rect(sl, x, y, w, 0.50, fill=color)
    tb(sl, f"Phase {num}", x + 0.12, y + 0.05, w - 0.24, 0.22, sz=10, bold=True, color=WHITE)
    tb(sl, title, x + 0.12, y + 0.26, w - 0.14, 0.22, sz=10, color=WHITE, bold=True)
    rect(sl, x, y + 0.50, w, h - 0.50, fill=LGRAY, line=color, lw=Pt(0.5))
    for i, it in enumerate(items):
        oval(sl, x + 0.14, y + 0.68 + i * 0.52, 0.09, 0.09, fill=color)
        tb(sl, it, x + 0.28, y + 0.62 + i * 0.52, w - 0.40, 0.48, sz=10, color=DARK)

def bullet_card(sl, title, items, x, y, w, h, hc=NAVY, bc=LGRAY, tc=DARK):
    rect(sl, x, y, w, 0.40, fill=hc)
    tb(sl, title, x + 0.14, y + 0.08, w - 0.20, 0.26, sz=11, bold=True, color=WHITE)
    rect(sl, x, y + 0.40, w, h - 0.40, fill=bc, line=hc, lw=Pt(0.5))
    for i, it in enumerate(items):
        oval(sl, x + 0.14, y + 0.58 + i * 0.48, 0.09, 0.09, fill=hc)
        tb(sl, it, x + 0.28, y + 0.52 + i * 0.48, w - 0.40, 0.44, sz=10, color=tc)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════
sl = ns()
rect(sl, 0, 0, 13.33, 3.60, fill=NAVY)
rect(sl, 0, 3.60, 13.33, 0.08, fill=AMBER)
# White title block
tb(sl, "AI-Powered Digital Transformation", 0.7, 0.55, 11.9, 0.85,
   sz=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(sl, "Intelligence Service", 0.7, 1.40, 11.9, 0.72,
   sz=36, bold=True, color=B8, align=PP_ALIGN.CENTER)
tb(sl, "A Consulting Service Offering for Enterprise Digital Transformation",
   1.0, 2.20, 11.33, 0.36, sz=14, color=B8,
   align=PP_ALIGN.CENTER, italic=True)
# Lower half — badges
for i, (lbl, col) in enumerate([
    ("TransformHub Platform", BLUE),
    ("18 AI LangGraph Agents", GREEN),
    ("5-Phase Methodology", PURPLE),
    ("Measurable Accuracy", AMBER),
]):
    bx = 1.0 + i * 2.90
    rect(sl, bx, 3.88, 2.55, 0.60, fill=col, line=col, lw=Pt(0))
    tb(sl, lbl, bx, 3.88, 2.55, 0.60, sz=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

tb(sl, "Powered by TransformHub  ·  Built for Enterprise-Scale Digital Transformation",
   1.0, 4.72, 11.33, 0.30, sz=10, color=MID, align=PP_ALIGN.CENTER, italic=True)

for i, (v, l) in enumerate([("$4.5T","DT Market 2027"),("73%","DT Failure Rate"),
                              ("95%","Faster Discovery"),("15×","Typical Client ROI")]):
    bx = 0.55 + i * 3.08
    metric_box(sl, v, l, bx, 5.22, w=2.72, h=1.00, vc=NAVY, bg=LGRAY, bc=MGRAY)

ftr(sl)
tp(sl, [
    "Open with: 'Every year $900B is wasted on digital transformations that fail — this offering changes that equation'",
    "TransformHub is the only AI-native platform with measurable, scored output accuracy per module",
    "This is not another consulting framework — it is an AI-powered intelligence layer that runs in days, not months",
    "The 15× ROI figure is based on: $7.6M value created vs ~$500K total engagement cost for a mid-size enterprise",
])

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Transformation Paradox (Hook)
# ══════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "The Transformation Paradox", "Why enterprises invest more — and deliver less — from digital transformation", NAVY)
ftr(sl)

# Big stat left
rect(sl, 0.40, CT, 3.90, CB - CT, fill=LGRAY, line=NAVY, lw=Pt(0.8))
tb(sl, "73%", 0.50, 1.40, 3.70, 1.20, sz=88, bold=True, color=RED, align=PP_ALIGN.CENTER)
tb(sl, "of digital transformations\nfail to meet objectives", 0.55, 2.70, 3.60, 0.60,
   sz=10, color=DARK, align=PP_ALIGN.CENTER)
tb(sl, "McKinsey Global Survey, 2024", 0.55, 3.42, 3.60, 0.24,
   sz=10, color=MID, italic=True, align=PP_ALIGN.CENTER)

# Right — the paradox statements
rect(sl, 4.65, CT, 8.30, CB - CT, fill=WHITE, line=MGRAY, lw=Pt(0.5))

stats = [
    ("$900B", "wasted annually on failed DT initiatives",          RED),
    ("3–5 yrs", "average time to see transformation value",         AMBER),
    ("6–8 tools", "siloed platforms used per enterprise",           BLUE),
    ("< 20%", "of capabilities mapped before transformation begins",PURPLE),
    ("$0",    "quantified accuracy score on AI-driven insights",    RED),
]
for i, (val, lbl, col) in enumerate(stats):
    y = CT + 0.10 + i * 0.76
    rect(sl, 4.75, y, 8.10, 0.68, fill=LGRAY, line=col, lw=Pt(1.2))
    rect(sl, 4.75, y, 0.14, 0.68, fill=col)
    tb(sl, val,  5.00, y + 0.08, 2.00, 0.36, sz=18, bold=True, color=col)
    tb(sl, lbl,  7.10, y + 0.18, 5.60, 0.32, sz=10, color=DARK)

tp(sl, [
    "Lead with empathy: 'Your leadership team believes transformation is on track — the data says otherwise'",
    "$900B in wasted DT spend is a Consortium for IT research figure — boards respond strongly to this number",
    "The '< 20% capabilities mapped' stat is the core gap we close — full AI-powered discovery in days",
    "Zero quantified accuracy is the status quo — every consulting recommendation is an educated guess today",
    "This paradox is WHY we built this service — not another framework, but measurable intelligence",
])

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Root Problem (6 pain points)
# ══════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Six Root Problems We Solve", "The structural gaps that doom digital transformation programmes", NAVY)
ftr(sl)

problems = [
    (RED,    "1", "Invisible Landscape",        "No systematic map of products, capabilities or functionalities exists before transformation begins"),
    (AMBER,  "2", "Unmapped Value Streams",     "Waste, bottlenecks and flow inefficiency are invisible — decisions made on gut feel, not data"),
    (BLUE,   "3", "Unquantified Risk",           "Compliance and operational risk is discovered mid-transformation — too late to avoid costly rework"),
    (PURPLE, "4", "Groundless Roadmaps",         "Transformation roadmaps built from workshops, not from AI-triangulated evidence and benchmarks"),
    (GREEN,  "5", "Tool Sprawl",                 "6–8 disconnected tools for process mining, risk, roadmap, architecture — no single source of truth"),
    (CYAN,   "6", "Black-Box AI",                "AI outputs have no confidence score — leaders can't distinguish high-quality insight from hallucination"),
]
for i, (col, num, title, desc) in enumerate(problems):
    col_i = i % 3; row_i = i // 3
    x = 0.40 + col_i * 4.26; y = CT + row_i * 1.96
    rect(sl, x, y, 4.06, 1.88, fill=LGRAY, line=col, lw=Pt(1.0))
    rect(sl, x, y, 4.06, 0.44, fill=col)
    tb(sl, f"#{num}  {title}", x + 0.14, y + 0.08, 3.78, 0.30, sz=11, bold=True, color=WHITE)
    tb(sl, desc, x + 0.14, y + 0.52, 3.78, 1.28, sz=10, color=DARK)

tp(sl, [
    "Problem #1 and #2 are the 'discovery gap' — together they make up 60% of wasted transformation spend",
    "Problem #3 (risk) is especially resonant in regulated industries: Banking, Insurance, Healthcare",
    "Problem #5 (tool sprawl) creates an immediate ROI story: one platform replacing 6–8 tools at $40–80K/year each",
    "Problem #6 (black-box AI) is our key differentiator — no competitor offers a per-module accuracy score",
])

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Why Traditional Approaches Fail
# ══════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Why Traditional Approaches Fail", "The consulting industry has not kept pace with AI-native transformation intelligence", NAVY)
ftr(sl)

rows = [
    ("Dimension",          "Traditional Consulting",    "TransformHub Approach",       True),
    ("Discovery Time",     "8–12 weeks of workshops",   "3–5 days, AI-automated",      False),
    ("Coverage",           "~20% of capabilities",      "100% in one pass",            False),
    ("Accuracy",           "Subjective / unscored",     "65–85% per-module, measured", False),
    ("VSM",                "Whiteboard sessions",       "AI-mapped with Mermaid diagrams", False),
    ("Risk Identification","Post-transformation",       "Pre-transformation with gates",False),
    ("Audit Trail",        "Word docs / spreadsheets",  "SHA-256 chained, tamper-proof",False),
    ("Cost",               "$400K–$800K typical",       "$250K–$350K typical",         False),
    ("Tool Count",         "6–8 separate platforms",    "1 integrated platform",       False),
]
rect(sl, 0.40, CT, 12.53, 3.82, fill=WHITE, line=MGRAY, lw=Pt(0.5))
for i, (dim, trad, ours, is_hdr) in enumerate(rows):
    y = CT + i * 0.44
    bg = NAVY if is_hdr else (LGRAY if i % 2 == 0 else WHITE)
    rect(sl, 0.40, y, 12.53, 0.42, fill=bg)
    c1 = WHITE if is_hdr else MID
    c2 = WHITE if is_hdr else RED
    c3 = WHITE if is_hdr else GREEN
    sz = 9.5 if is_hdr else 9
    tb(sl, dim,  0.55, y + 0.10, 2.80, 0.26, sz=sz, bold=is_hdr, color=c1)
    tb(sl, trad, 3.50, y + 0.10, 4.40, 0.26, sz=sz, color=c2)
    tb(sl, ours, 8.05, y + 0.10, 4.70, 0.26, sz=sz, bold=(not is_hdr), color=c3)

tp(sl, [
    "Use this slide to 'de-position' traditional discovery: 12 weeks of workshops for 20% coverage is no longer acceptable",
    "The audit trail row lands hardest in regulated industries — SHA-256 chaining is a compliance differentiator",
    "Cost comparison is conservative — many firms see $600K–$1M in traditional consulting for the same scope",
    "Emphasise that our accuracy score is unique — no competitor in this space offers measurable AI output quality",
])

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — The Market Opportunity
# ══════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "The Market Opportunity", "A $4.5 trillion market looking for better answers", NAVY)
ftr(sl)

for i, (v, l, sub, col, bg) in enumerate([
    ("$4.5T",  "Global DT Market by 2027",  "CAGR 23% (IDC 2024)",         NAVY,   LGRAY),
    ("$180B",  "Process Intelligence Spend", "CAGR 31% (Gartner)",          BLUE,   BLUEBG),
    ("$42B",   "AI in Enterprise Consulting","Fastest growing segment",      GREEN,  GREENBG),
    ("73%",    "Failure Rate",               "Creates massive demand for better", RED, REDBG),
]):
    bx = 0.40 + i * 3.16
    metric_box(sl, v, l, bx, CT, w=2.96, h=1.10, vc=col, bg=bg, bc=col)
    tb(sl, sub, bx, CT + 1.14, 2.96, 0.24, sz=10, color=MID, align=PP_ALIGN.CENTER, italic=True)

# Target segments
tb(sl, "Target Segments & Entry Points", 0.40, CT + 1.58, 12.0, 0.32, sz=10, bold=True, color=NAVY)
segs = [
    (BLUE,   "Banking & Financial Services", ["SOX / FDIC / GLBA compliance pressure", "Digital banking capability gaps", "Payments transformation at scale"]),
    (GREEN,  "Healthcare & Life Sciences",   ["HIPAA / clinical workflow modernisation", "EHR integration complexity", "Patient journey digitisation"]),
    (PURPLE, "Insurance",                    ["Claims process automation", "Underwriting AI readiness", "RegTech compliance mapping"]),
    (AMBER,  "Retail & eCommerce",           ["Omnichannel value stream gaps", "Supply chain digitisation", "Customer experience transformation"]),
    (CYAN,   "Telco & Media",                ["Network capability modernisation", "5G product portfolio mapping", "BSS/OSS transformation"]),
]
for i, (col, name, pts) in enumerate(segs):
    x = 0.40 + i * 2.52
    rect(sl, x, CT + 1.94, 2.42, 2.82, fill=LGRAY, line=col, lw=Pt(0.8))
    rect(sl, x, CT + 1.94, 2.42, 0.38, fill=col)
    tb(sl, name, x + 0.10, CT + 1.98, 2.22, 0.30, sz=10, bold=True, color=WHITE)
    for j, p in enumerate(pts):
        oval(sl, x + 0.12, CT + 2.46 + j * 0.68, 0.09, 0.09, fill=col)
        tb(sl, p, x + 0.26, CT + 2.38 + j * 0.68, 2.10, 0.62, sz=10, color=DARK)

tp(sl, [
    "Banking and Financial Services is the highest-value entry: regulatory pressure creates immediate urgency and budget",
    "Healthcare is fastest-growing: post-pandemic digital mandates + HIPAA complexity = strong ROI argument",
    "Lead with a single business unit POC — land and expand is the proven GTM motion for this type of platform",
    "The $42B AI in enterprise consulting market is our direct addressable market — this is the fastest-growing segment",
])

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Service Offering Overview
# ══════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Service Offering: TransformHub Intelligence Service", "End-to-end AI-powered transformation intelligence — from discovery to execution roadmap", NAVY)
ftr(sl)

# Central spine
rect(sl, 0.40, CT, 12.53, 0.52, fill=NAVY)
tb(sl, "THE 5-PHASE TRANSFORMATION INTELLIGENCE METHODOLOGY", 0.60, CT + 0.12,
   12.0, 0.30, sz=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

phases = [
    (RED,    "01 DISCOVER",   "Digital Landscape\nAssessment"),
    (AMBER,  "02 MAP",        "Value Stream\nAnalysis"),
    (BLUE,   "03 DESIGN",     "Future State\nVision"),
    (PURPLE, "04 GOVERN",     "Risk &\nCompliance"),
    (GREEN,  "05 EXECUTE",    "Roadmap &\nEnablement"),
]
for i, (col, num, name) in enumerate(phases):
    x = 0.40 + i * 2.52
    rect(sl, x, CT + 0.56, 2.42, 0.80, fill=col)
    tb(sl, num,  x + 0.10, CT + 0.60, 2.22, 0.28, sz=10,  bold=True, color=WHITE)
    tb(sl, name, x + 0.10, CT + 0.88, 2.22, 0.44, sz=10, color=WHITE)
    if i < 4:
        tb(sl, "→", x + 2.42, CT + 0.80, 0.20, 0.28, sz=14, bold=True, color=MID, align=PP_ALIGN.CENTER)

# Key stats below phases
for i, (val, lbl, col) in enumerate([
    ("3–5 days", "Full discovery",       RED),
    ("100%",     "VSM coverage",         AMBER),
    ("3-band",   "Projected metrics",    BLUE),
    ("SHA-256",  "Audit trail",          PURPLE),
    ("25+",      "Deliverables",         GREEN),
]):
    x = 0.40 + i * 2.52
    rect(sl, x, CT + 1.46, 2.42, 0.76, fill=LGRAY, line=col, lw=Pt(0.6))
    tb(sl, val, x, CT + 1.52, 2.42, 0.36, sz=16, bold=True, color=col, align=PP_ALIGN.CENTER)
    tb(sl, lbl, x, CT + 1.88, 2.42, 0.28, sz=10, color=MID, align=PP_ALIGN.CENTER)

# What it includes row
tb(sl, "What This Offering Includes:", 0.40, CT + 2.36, 4.0, 0.28, sz=10, bold=True, color=NAVY)
includes = [
    ("🤖", "18 AI LangGraph Agents   running continuously"),
    ("📊", "Per-module accuracy scoring  on every output"),
    ("🔗", "RAG pipeline grounded in   client's own documents"),
    ("🛡️", "Risk gates + compliance   framework mapping"),
    ("🗺️", "End-to-end Mermaid VSM   diagrams for all capabilities"),
    ("📋", "Executive-ready deliverables  at every phase gate"),
]
for i, (icon, txt) in enumerate(includes):
    col_i = i % 3; row_i = i // 3
    x = 0.40 + col_i * 4.18; y = CT + 2.70 + row_i * 0.68
    rect(sl, x, y, 4.00, 0.60, fill=LGRAY, line=MGRAY, lw=Pt(0.5))
    tb(sl, icon, x + 0.10, y + 0.10, 0.40, 0.40, sz=14)
    tb(sl, txt,  x + 0.58, y + 0.10, 3.30, 0.44, sz=10, color=DARK)

tp(sl, [
    "This offering wraps the TransformHub platform in a repeatable 5-phase consulting engagement",
    "AI agents do the heavy lifting — our consultants guide, interpret and present the outputs",
    "Emphasise the RAG pipeline: AI outputs are grounded in the client's own uploaded documents — not generic models",
    "The accuracy scoring capability is a breakthrough — clients can show their board a % score for AI output quality",
    "25+ deliverables means every phase has a tangible artefact the client owns after the engagement",
])

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Objectives
# ══════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Service Objectives", "What this offering is designed to achieve for every client", NAVY)
ftr(sl)

objectives = [
    (RED,    "Strategic Objective 1",  "Compress Discovery",
     "Reduce the time to build a complete digital landscape map from 8–12 weeks to 3–5 days using AI-powered multi-pass discovery across all repositories, APIs, database schemas and documentation.",
     ["3–5 day discovery (vs 8–12 weeks)", "100% capability coverage (vs ~20%)", "Confidence scores on every entity", "8 evidence source types triangulated"]),
    (BLUE,   "Strategic Objective 2",  "Ground Every Decision in Data",
     "Replace gut-feel roadmaps with AI-triangulated evidence: VSM data grounded in uploaded benchmarks, projected metrics with conservative/expected/optimistic bands, and RICE-scored prioritisation.",
     ["Benchmark-grounded future state", "3-band projected metrics", "RICE-scored transformation priorities", "BM25 + vector hybrid RAG retrieval"]),
    (GREEN,  "Strategic Objective 3",  "Deliver Measurable, Auditable Intelligence",
     "Provide the first transformation advisory service with a quantifiable accuracy score — per-module, per-org — so leadership can evaluate AI output quality and improve it over time.",
     ["65–85% composite accuracy score", "SHA-256 audit trail (tamper-proof)", "Human-in-the-loop approval gates", "Improving scores as agents learn"]),
]
for i, (col, eyebrow, title, desc, bullets) in enumerate(objectives):
    y = CT + i * 1.30
    rect(sl, 0.40, y, 12.53, 1.22, fill=LGRAY, line=col, lw=Pt(1.0))
    rect(sl, 0.40, y, 0.18, 1.22, fill=col)
    tb(sl, eyebrow, 0.70, y + 0.06, 3.00, 0.22, sz=10, bold=True, color=col)
    tb(sl, title,   0.70, y + 0.28, 3.20, 0.36, sz=14, bold=True, color=DARK)
    tb(sl, desc,    0.70, y + 0.66, 5.80, 0.52, sz=10, color=MID)
    for j, b in enumerate(bullets):
        oval(sl, 6.80, y + 0.22 + j * 0.26, 0.09, 0.09, fill=col)
        tb(sl, b, 6.96, y + 0.16 + j * 0.26, 5.80, 0.26, sz=10, color=DARK, bold=True)

tp(sl, [
    "Objective 1 creates the 'wow moment' — show them the Discovery page with 78 capabilities already mapped",
    "Objective 2 directly addresses the board: 'Our roadmap is grounded in your own uploaded benchmarks, not our opinion'",
    "Objective 3 is our moat: accuracy scoring is unique in the market — no big-4 firm or software vendor offers this",
    "These 3 objectives map directly to the 3 questions every CIO asks: What do we have? Where should we go? How sure are you?",
])

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Quantitative Benefits
# ══════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Quantitative Benefits", "Hard numbers from pilot engagements and platform benchmarks", NAVY)
ftr(sl)

tb(sl, "SPEED & EFFICIENCY GAINS", 0.40, CT, 6.00, 0.28, sz=10, bold=True, color=NAVY)
speed = [
    ("Discovery cycle time",       "8–12 weeks",  "3–5 days",   "95%↓",  GREEN),
    ("VSM capability coverage",    "~20%",        "100%",       "5×↑",   GREEN),
    ("Time to transformation roadmap","3–4 months","2–3 weeks", "85%↓",  GREEN),
    ("Analyst productivity (products/person)","1–2","12–15",    "8×↑",   GREEN),
    ("Tool consolidation",         "6–8 tools",   "1 platform", "75%↓",  GREEN),
]
rect(sl, 0.40, CT + 0.32, 6.20, 3.70, fill=WHITE, line=MGRAY, lw=Pt(0.5))
rect(sl, 0.40, CT + 0.32, 6.20, 0.34, fill=NAVY)
for j, hd in enumerate(["Metric","Before","After","Gain"]):
    tb(sl, hd, 0.55 + j * 1.52, CT + 0.38, 1.45, 0.22, sz=10, bold=True, color=WHITE)
for i, (m, b, a, g, c) in enumerate(speed):
    y = CT + 0.70 + i * 0.64
    bg = LGRAY if i % 2 == 0 else WHITE
    rect(sl, 0.40, y, 6.20, 0.62, fill=bg)
    tb(sl, m, 0.55, y + 0.18, 2.00, 0.26, sz=10, color=DARK)
    tb(sl, b, 2.62, y + 0.18, 1.42, 0.26, sz=10, color=MID)
    tb(sl, a, 4.10, y + 0.18, 1.42, 0.26, sz=10, color=B4, bold=True)
    tb(sl, g, 5.55, y + 0.18, 1.00, 0.26, sz=10, bold=True, color=c)

tb(sl, "FINANCIAL & QUALITY IMPACT", 6.80, CT, 6.0, 0.28, sz=10, bold=True, color=NAVY)
fin = [
    ("$200K–400K",  "Cost saved vs traditional consulting per engagement",      BLUE),
    ("$300K–500K",  "Annual savings from tool consolidation",                   GREEN),
    ("65–85%",      "AI output accuracy score (from zero/unscored baseline)",   NAVY),
    ("3–5×",        "ROI for client in first 18 months",                        GREEN),
    ("2–3 CRITICAL","Risks identified and mitigated pre-transformation",        RED),
    ("6 months",    "Faster time-to-market for priority digital products",      BLUE),
]
for i, (val, lbl, col) in enumerate(fin):
    col_i = i % 2; row_i = i // 2
    x = 6.90 + col_i * 3.08; y = CT + 0.36 + row_i * 0.90
    rect(sl, x, y, 2.90, 0.82, fill=LGRAY, line=col, lw=Pt(0.8))
    tb(sl, val, x, y + 0.04, 2.90, 0.42, sz=20, bold=True, color=col, align=PP_ALIGN.CENTER)
    tb(sl, lbl, x + 0.10, y + 0.46, 2.70, 0.32, sz=10, color=MID, align=PP_ALIGN.CENTER)

tp(sl, [
    "The 95% reduction in discovery time is the single most powerful number — anchor the conversation here",
    "Analyst productivity 8× improvement means one consultant can now serve 12–15 products vs 1–2 traditionally",
    "Tool consolidation ROI is easily calculable with the client: ask 'how many tools do you use today, and what's the licence cost?'",
    "The 65–85% accuracy score is an industry first — emphasise that before our platform, this number was literally unmeasurable",
])

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Why Critical Now
# ══════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Why This is Critical to Implement Now", "Five forces making 2025–2026 the decisive window for AI-native transformation", NAVY)
ftr(sl)

forces = [
    (RED,    "⚖️  Regulatory Pressure",
     "EU AI Act (Aug 2026), DORA (Jan 2025), SEC AI disclosure rules and APRA CPS 230 all require documented, auditable AI decision trails. Organisations without an audit-capable platform face fines and remediation costs of $5M–$50M.",
     "Fines up to €30M or 6% of global revenue under EU AI Act"),
    (AMBER,  "🏃  Competitor Velocity",
     "67% of financial services firms have AI transformation programmes underway (Accenture 2024). The window to establish AI-native capabilities before competitors embed them is 18–24 months. After that, the gap becomes structural.",
     "First-mover advantage window: 18–24 months remaining"),
    (BLUE,   "💸  Technical Debt Crisis",
     "Legacy technical debt is accelerating: the average enterprise spends 72% of IT budget maintaining existing systems (Gartner). Without a systematic capability map, organisations cannot identify which systems to modernise first.",
     "$3.61T in global technical debt — growing 15% per year"),
    (PURPLE, "📊  Board Accountability Shift",
     "87% of boards now require quantified progress metrics for digital transformation programmes (KPMG 2024). Vague 'transformation journeys' are no longer acceptable — executives need scored, auditable output.",
     "87% of boards now demand measurable DT metrics"),
    (GREEN,  "🤖  AI Commoditisation Window",
     "AI-native consulting tools are 2–3 years from commoditisation. The organisations that establish AI-augmented transformation capabilities now will have 3–5 years of compounding advantage before the market equalises.",
     "AI consulting commoditisation: estimated 2027–2028"),
]
for i, (col, title, desc, stat) in enumerate(forces):
    col_i = i % 3 if i < 3 else (i - 3)
    row_i = 0 if i < 3 else 1
    if i < 3:
        x = 0.40 + i * 4.22; y = CT
        w = 4.02; h = 2.10
    else:
        x = 0.40 + (i-3) * 6.38; y = CT + 2.18
        w = 6.18; h = 1.64
    rect(sl, x, y, w, 0.38, fill=col)
    tb(sl, title, x + 0.12, y + 0.07, w - 0.20, 0.26, sz=10, bold=True, color=WHITE)
    rect(sl, x, y + 0.38, w, h - 0.38, fill=LGRAY, line=col, lw=Pt(0.5))
    tb(sl, desc, x + 0.14, y + 0.46, w - 0.22, h - 0.68, sz=10, color=DARK)
    tb(sl, stat, x + 0.14, y + h - 0.28, w - 0.22, 0.22, sz=10, bold=True, color=col, italic=True)

tp(sl, [
    "Regulatory slide is the opener for banking/insurance clients — DORA enforcement began Jan 2025, AI Act by Aug 2026",
    "Competitor velocity: 'Your peers are already moving — the question is whether you move with data or without it'",
    "Technical debt stat ($3.61T) reframes the value: this isn't a nice-to-have, it's existential budget allocation",
    "Board accountability shift is the internal sponsor's best argument for getting budget approved",
])

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Consulting Approach Overview
# ══════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Consulting Methodology: 5-Phase Approach", "A structured, repeatable engagement model with AI-augmented delivery at every phase", NAVY)
ftr(sl)

phase_data = [
    (RED,    "01", "DISCOVER",  "Weeks 1–2",  ["Org onboarding","AI repo/API analysis","3-pass capability discovery","Stakeholder interviews","Confidence scoring"]),
    (AMBER,  "02", "MAP",       "Weeks 3–4",  ["L1/L2/L3 VSM build","Bottleneck identification","Flow efficiency baseline","Mermaid diagram generation","Benchmark comparison"]),
    (BLUE,   "03", "DESIGN",    "Weeks 5–6",  ["Automation mix analysis","Future state architecture","3-band projection model","RICE prioritisation","Tech stack assessment"]),
    (PURPLE, "04", "GOVERN",    "Week 7",     ["Risk identification (AI)","Compliance framework map","SHA-256 audit setup","Risk gate thresholds","Mitigation planning"]),
    (GREEN,  "05", "EXECUTE",   "Week 8",     ["Roadmap construction","Architecture blueprint","Exec presentation","KB population","Accuracy baseline"]),
]
for i, (col, num, phase, timing, items) in enumerate(phase_data):
    x = 0.40 + i * 2.52
    phase_card(sl, num + " " + phase, timing, col, items, x, CT, w=2.38, h=3.62)

# Gate labels between phases
for i in range(4):
    x = 2.58 + i * 2.52
    rect(sl, x + 0.02, CT + 1.60, 0.22, 0.44, fill=LGRAY, line=MGRAY, lw=Pt(0.5))
    tb(sl, "✓\nGate", x + 0.04, CT + 1.62, 0.18, 0.40, sz=7, color=MID, align=PP_ALIGN.CENTER)

# Deliverables count bar at bottom
rect(sl, 0.40, CT + 3.70, 12.53, 0.46, fill=NAVY)
tb(sl, "Total Deliverables:  25+  artefacts  ·  Every phase has a client-owned, executive-ready output",
   0.55, CT + 3.80, 12.20, 0.28, sz=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

tp(sl, [
    "Each phase ends with a gate review — client approves before the next phase begins (HITL = human-in-the-loop)",
    "Phases 1–4 are typically 8 weeks; Phase 5 can extend to 12 weeks for complex enterprise environments",
    "The AI agents run throughout all phases — consultants guide, interpret and present, not do manual analysis",
    "Gate checkpoints are a key selling point for clients: no cost overrun risk, clear go/no-go at each milestone",
])

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Phase 1: Discover
# ══════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Phase 1: Digital Landscape Discovery", "Weeks 1–2  ·  AI-powered, multi-pass discovery of all digital products, capabilities and functionalities", RED)
ftr(sl)

rect(sl, 0.40, CT, 0.24, CB - CT, fill=RED)
# Activities
bullet_card(sl, "Key Activities", [
    "Day 1–2: Organisation onboarding — connect to GitHub, OpenAPI specs, DB schemas, test suites",
    "Day 2–3: AI Pass 1 — identify all Digital Products with human review gate",
    "Day 3–4: AI Pass 2 — map Capabilities per approved product with confidence scores",
    "Day 4–5: AI Pass 3 — extract Functionalities with source attribution (8 evidence types)",
    "Day 5–8: Stakeholder workshops to validate AI findings (not to gather them)",
    "Day 9–10: Business segment classification + repository mapping",
], 0.76, CT, 5.60, CB - CT, hc=RED, bc=LGRAY)

# Deliverables
bullet_card(sl, "Deliverables", [
    "D1.1  Digital Product Inventory (9–15 products typical)",
    "D1.2  Capability Map with confidence scores (78–120 capabilities)",
    "D1.3  Functionality Decomposition (200–400 functionalities)",
    "D1.4  Source Attribution Report (GitHub, OpenAPI, DB, Tests)",
    "D1.5  Discovery Accuracy Baseline Report (~85% typical)",
    "D1.6  Executive Summary — landscape overview (board-ready)",
], 6.52, CT, 6.42, CB - CT, hc=RED, bc=LGRAY)

# Stats bar
rect(sl, 0.40, CB - 0.02, 12.53, 0.56, fill=RGBColor(0xFE, 0xF2, 0xF2), line=RED, lw=Pt(0.5))
for i, (v, l) in enumerate([("3–5 days","Full discovery")  ,("95%","Faster than workshops"),
                              ("8 sources","Evidence triangulation"),("85%","Avg accuracy score")]):
    bx = 0.70 + i * 3.0
    tb(sl, v, bx, CB - 0.01, 2.80, 0.30, sz=13, bold=True, color=RED, align=PP_ALIGN.CENTER)
    tb(sl, l, bx, CB + 0.28, 2.80, 0.20, sz=10, color=MID, align=PP_ALIGN.CENTER)

tp(sl, [
    "The key message: AI does the analysis, consultants do the interpretation — discovery is 95% faster than workshops",
    "Confidence scoring is the trust-builder: every entity shows 78–97% confidence + 3–5 source types used",
    "Human review gate after Pass 1 means client approves the product list before AI maps capabilities — no surprises",
    "D1.6 Executive Summary is the first board-ready artefact — visible value by end of week 1",
])

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Phase 2: VSM
# ══════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Phase 2: Value Stream Analysis", "Weeks 3–4  ·  Three-level VSM across all capabilities — bottlenecks, waste and flow efficiency quantified", AMBER)
ftr(sl)

rect(sl, 0.40, CT, 0.24, CB - CT, fill=AMBER)
bullet_card(sl, "Key Activities", [
    "L1 Analysis: Segment-level flow across all products — identify which products are bottlenecked",
    "L2 Analysis: Product-level capability-by-capability metrics (process time, wait time, flow efficiency)",
    "L3 Analysis: Functionality-level step classification — Value-Adding / Bottleneck / Waste / Waiting",
    "Mermaid diagram auto-generation for every capability (100% coverage, not sampled)",
    "Industry benchmark comparison using uploaded VSM KPI benchmarks",
    "Bottleneck register creation — rank waste by business impact",
], 0.76, CT, 5.60, CB - CT, hc=AMBER, bc=AMBERBG)

bullet_card(sl, "Deliverables", [
    "D2.1  L1/L2/L3 VSM Diagrams — Mermaid visualisations for all capabilities",
    "D2.2  Flow Efficiency Baseline Report (current state, per capability)",
    "D2.3  Bottleneck & Waste Register (ranked by business impact)",
    "D2.4  Industry Benchmark Comparison (process time, cycle time, FE%)",
    "D2.5  VSM Accuracy Report (~72% typical for seeded, higher for live)",
    "D2.6  L1 Executive Overview — cross-product flow heat map",
], 6.52, CT, 6.42, CB - CT, hc=AMBER, bc=AMBERBG)

rect(sl, 0.40, CB - 0.02, 12.53, 0.56, fill=AMBERBG, line=AMBER, lw=Pt(0.5))
for i, (v, l) in enumerate([("100%","VSM capability coverage"),("78+","Mermaid diagrams generated"),
                              ("~35%","Typical baseline flow efficiency"),("4 types","Step classification")]):
    bx = 0.70 + i * 3.0
    tb(sl, v, bx, CB - 0.01, 2.80, 0.30, sz=13, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
    tb(sl, l, bx, CB + 0.28, 2.80, 0.20, sz=10, color=MID, align=PP_ALIGN.CENTER)

tp(sl, [
    "100% VSM coverage is the key stat — traditional VSM workshops sample 10–20% at best, over multiple weeks",
    "Mermaid diagrams are client-owned, version-controlled artefacts — they can embed them in Confluence/SharePoint",
    "Flow efficiency of ~35% baseline is typical — the 'waste' visible in the data is where business value hides",
    "L3 step classification directly feeds the Future State Design phase — no re-analysis needed, it flows forward",
])

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Phase 3: Design
# ══════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Phase 3: Future State Vision & Design", "Weeks 5–6  ·  Benchmark-grounded transformation design with automation mix, 3-band projections and RICE priorities", BLUE)
ftr(sl)

rect(sl, 0.40, CT, 0.24, CB - CT, fill=BLUE)
bullet_card(sl, "Key Activities", [
    "Upload VSM benchmark documents & transformation case studies to Context Hub (RAG grounding)",
    "Run Future State Vision agent — AI generates capability modernisation roadmap per product",
    "Automation mix analysis: RPA / AI-ML / Agent-Based / Conversational / Analytics breakdown",
    "3-band projection modelling: Conservative / Expected / Optimistic per capability",
    "RICE prioritisation scoring: Reach × Impact × Confidence ÷ Effort for each initiative",
    "Future state architecture design: target tech stack, integration patterns, cloud alignment",
], 0.76, CT, 5.60, CB - CT, hc=BLUE, bc=BLUEBG)

bullet_card(sl, "Deliverables", [
    "D3.1  Future State Vision Report — per-product transformation narrative",
    "D3.2  Automation Mix Breakdown — RPA/AI/Agent proportions per product",
    "D3.3  Projected Metrics Report — 3-band forecasts (Conservative/Expected/Optimistic)",
    "D3.4  Capability Modernisation Register — RICE-scored priority list",
    "D3.5  Target Architecture Blueprint — initial draft (refined in Phase 5)",
    "D3.6  Benchmark Comparison Appendix — AI projections vs industry data",
], 6.52, CT, 6.42, CB - CT, hc=BLUE, bc=BLUEBG)

rect(sl, 0.40, CB - 0.02, 12.53, 0.56, fill=BLUEBG, line=BLUE, lw=Pt(0.5))
for i, (v, l) in enumerate([("3-band","Projection scenarios"),("Benchmark","Grounded via RAG"),
                              ("RICE","Scored priorities"),("5 types","Automation mix analysis")]):
    bx = 0.70 + i * 3.0
    tb(sl, v, bx, CB - 0.01, 2.80, 0.30, sz=13, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    tb(sl, l, bx, CB + 0.28, 2.80, 0.20, sz=10, color=MID, align=PP_ALIGN.CENTER)

tp(sl, [
    "The 'benchmark-grounded' differentiator: AI reads the client's own uploaded case studies to generate projections",
    "3-band projections give clients a credible range rather than a single point estimate — boards trust ranges more",
    "RICE scoring removes politics from prioritisation — data-driven ranking, not HiPPO-driven (Highest Paid Person's Opinion)",
    "Phase 3 output is often the most impressive demo moment: future state with automation mix charts and projected flow efficiency",
])

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Phase 4: Govern
# ══════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Phase 4: Risk & Compliance Governance", "Week 7  ·  AI-powered risk identification, compliance framework mapping and tamper-proof audit trail", PURPLE)
ftr(sl)

rect(sl, 0.40, CT, 0.24, CB - CT, fill=PURPLE)
bullet_card(sl, "Key Activities", [
    "AI-assisted risk identification across all capabilities and transformation initiatives",
    "Risk scoring: CRITICAL (≥8.0) / HIGH (6.0–7.9) / MEDIUM / LOW with auto gate assignment",
    "Compliance framework mapping: SOX, FDIC, GLBA, GDPR, DORA, HIPAA, PCI-DSS (configurable)",
    "SHA-256 chained audit trail setup — every state change cryptographically recorded",
    "Risk gate configuration: CRITICAL risks automatically block transformation progression",
    "Mitigation plan development with owner, timeline and evidence link requirements",
], 0.76, CT, 5.60, CB - CT, hc=PURPLE, bc=RGBColor(0xF5, 0xF3, 0xFF))

bullet_card(sl, "Deliverables", [
    "D4.1  Risk Assessment Register (by category: Compliance, Security, Operational, Performance)",
    "D4.2  Compliance Framework Mapping (applicable regulations per capability)",
    "D4.3  SHA-256 Audit Trail Report (cryptographic proof of all changes)",
    "D4.4  Risk Gate Configuration Document (thresholds and escalation paths)",
    "D4.5  Mitigation Action Plan (owner, timeline, evidence for each CRITICAL/HIGH risk)",
    "D4.6  Regulatory Readiness Score — % of compliance requirements met",
], 6.52, CT, 6.42, CB - CT, hc=PURPLE, bc=RGBColor(0xF5, 0xF3, 0xFF))

rect(sl, 0.40, CB - 0.02, 12.53, 0.56, fill=RGBColor(0xF5, 0xF3, 0xFF), line=PURPLE, lw=Pt(0.5))
for i, (v, l) in enumerate([("SHA-256","Chained audit trail"),("Auto-block","CRITICAL risk gate"),
                              ("10+","Compliance frameworks"),("Pre-transform","Risk identification")]):
    bx = 0.70 + i * 3.0
    tb(sl, v, bx, CB - 0.01, 2.80, 0.30, sz=13, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    tb(sl, l, bx, CB + 0.28, 2.80, 0.20, sz=10, color=MID, align=PP_ALIGN.CENTER)

tp(sl, [
    "SHA-256 chained audit trail is a must-have for banking, insurance and healthcare — DORA mandates it from Jan 2025",
    "CRITICAL risk auto-block is a governance breakthrough: no human can inadvertently proceed past a critical risk",
    "Pre-transformation risk identification is the key message: finding risks in Phase 4 costs 10× less than finding them in production",
    "The Regulatory Readiness Score gives the CISO and CRO a single number to track — and it improves as risks are mitigated",
])

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Phase 5: Execute
# ══════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Phase 5: Transformation Roadmap & Enablement", "Week 8  ·  Transformation roadmap, architecture blueprint, knowledge base and client team enablement", GREEN)
ftr(sl)

rect(sl, 0.40, CT, 0.24, CB - CT, fill=GREEN)
bullet_card(sl, "Key Activities", [
    "Transformation roadmap construction: RICE-scored initiatives in quarterly delivery tracks",
    "Architecture blueprint: target state, integration patterns, cloud platform selection",
    "Executive presentation package: board-ready visuals for each phase's key findings",
    "Context Hub population: seed knowledge base with client's key documents (BRDs, specs, policies)",
    "Accuracy scoring and baseline: per-module composite score for the full engagement",
    "Client team enablement: 2-day platform training for internal transformation team",
], 0.76, CT, 5.60, CB - CT, hc=GREEN, bc=GREENBG)

bullet_card(sl, "Deliverables", [
    "D5.1  Transformation Roadmap — 4-quarter delivery plan, RICE-ranked",
    "D5.2  Architecture Blueprint — target state + integration architecture",
    "D5.3  Executive Presentation Deck — all phases, board-ready",
    "D5.4  Accuracy Baseline Report — per-module scores with improvement plan",
    "D5.5  Knowledge Base Package — seeded context docs for ongoing AI grounding",
    "D5.6  Platform Handover Guide — for internal team to continue using TransformHub",
], 6.52, CT, 6.42, CB - CT, hc=GREEN, bc=GREENBG)

rect(sl, 0.40, CB - 0.02, 12.53, 0.56, fill=GREENBG, line=GREEN, lw=Pt(0.5))
for i, (v, l) in enumerate([("25+","Total deliverables"),("4-quarter","Roadmap tracks"),
                              ("Board-ready","Executive deck"),("Ongoing","Platform access")]):
    bx = 0.70 + i * 3.0
    tb(sl, v, bx, CB - 0.01, 2.80, 0.30, sz=13, bold=True, color=B4, align=PP_ALIGN.CENTER)
    tb(sl, l, bx, CB + 0.28, 2.80, 0.20, sz=10, color=MID, align=PP_ALIGN.CENTER)

tp(sl, [
    "Phase 5 is where consulting value is most visible: the executive deck lets the partner present findings to the board directly",
    "Platform handover is a key differentiator: client keeps TransformHub running with their own team after the engagement ends",
    "The accuracy baseline report creates an upsell opportunity: 'You're at 68% today — run the Risk Agent to get to 75%'",
    "Knowledge base seeding is the most important long-term investment: better RAG grounding means better AI outputs forever",
])

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Illustrative Deliverables
# ══════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Illustrative Deliverables", "Representative samples of the 25+ client-owned artefacts produced across the 5 phases", NAVY)
ftr(sl)

deliverables = [
    (RED,    "D1 Discovery Accuracy Report",
     "Executive-ready PDF showing 9 products, 26 capabilities, 79 functionalities discovered for US Bank's Retail Banking segment. Includes per-entity confidence scores (78–97%), 5 evidence source types used, and an accuracy composite of 85%."),
    (AMBER,  "D2 VSM Flow Efficiency Dashboard",
     "Interactive L1→L2→L3 drill-down showing 100% capability coverage. LoanFlow Digital: 72% flow efficiency, 4hr process time, 2.2hr wait time. LegacyCore: 23% flow efficiency — CRITICAL bottleneck flagged. Mermaid diagrams attached for all 26 capabilities."),
    (BLUE,   "D3 Future State Vision Report",
     "Per-product transformation narrative grounded in uploaded Gartner VSM benchmarks. LoanFlow target: 91% flow efficiency (+19pts). Automation mix: 45% AI/ML, 30% RPA, 15% Agent-Based, 10% Conversational. Conservative/Expected/Optimistic metrics included."),
    (PURPLE, "D4 Risk Gate & Compliance Report",
     "4 risk assessments, 2 CRITICAL risks blocking transformation (SHA-256 hash: a3f7b9c2...). Compliance mapping: 6 frameworks (SOX, PCI-DSS, FDIC, GLBA, FFIEC, NIST). 4 chained audit entries. Regulatory Readiness Score: 74%."),
    (GREEN,  "D5 Transformation Roadmap",
     "4-quarter delivery roadmap: Q1 LoanFlow AI modernisation (RICE: 38), Q2 CoreBanking API layer (RICE: 31), Q3 Customer Data Platform (RICE: 28), Q4 Analytics Hub (RICE: 22). Total initiative portfolio: $12.4M budget, 18-month ROI: $37M."),
    (CYAN,   "D6 Accuracy Baseline & Action Plan",
     "Composite score: 68% (Discovery 85%, VSM 72%, Future State 70%, Risk 35%). Action Plan: 6 prioritised improvement steps — Run Risk Agent for 6 more products (+8pts), Upload architecture standards doc (+4pts), Run PT agent for all products (+5pts)."),
]
for i, (col, title, desc) in enumerate(deliverables):
    col_i = i % 3; row_i = i // 3
    x = 0.40 + col_i * 4.22; y = CT + row_i * 1.94
    rect(sl, x, y, 4.02, 1.84, fill=LGRAY, line=col, lw=Pt(1.0))
    rect(sl, x, y, 4.02, 0.34, fill=col)
    tb(sl, title, x + 0.12, y + 0.06, 3.78, 0.24, sz=10, bold=True, color=WHITE)
    tb(sl, desc,  x + 0.12, y + 0.42, 3.78, 1.36, sz=10, color=DARK)

tp(sl, [
    "Walk through D1 and D2 during the demo — they are the most viscerally impressive to non-technical stakeholders",
    "D4 (Risk Gate Report) with SHA-256 hash values is what gets the attention of CROs, CISOs and Compliance teams",
    "D5 (Transformation Roadmap) is the client's primary take-away — it should look polished enough for a board pack",
    "D6 (Accuracy Baseline) creates the 'continuous improvement' narrative and often drives upsell into follow-on engagements",
])

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Critical Success Factors
# ══════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Critical Success Factors", "What needs to be true for this offering to deliver its full value — and what we control vs the client controls", NAVY)
ftr(sl)

csfs_ours = [
    ("Consultant AI Literacy",     "Delivery team must understand how LangGraph agents work, how to interpret confidence scores and how to QA agent outputs — not just present them"),
    ("RAG Quality Management",     "Context Hub must be seeded with high-quality benchmark documents before agents run — garbage in, garbage out applies doubly to RAG pipelines"),
    ("Engagement Management",      "A dedicated EM who manages client access, gate approvals and stakeholder scheduling — under-managed engagements lose 40% of value"),
    ("Interpretation Layer",       "Raw AI outputs must be interpreted, contextualised and presented by consultants — clients pay for insight, not raw data"),
]
csfs_client = [
    ("C-Suite Sponsorship",        "The engagement must have a named executive sponsor at CXO level who champions access, budget and organisational buy-in"),
    ("Source System Access",       "Access to GitHub repos, OpenAPI specs, DB schemas and test suites is required within 48 hours of kick-off — delays cascade"),
    ("Domain SME Availability",    "2–3 hours/week from 1–2 business domain SMEs for Phase 1 and 3 gate reviews — non-negotiable for accuracy"),
    ("Clear Business Segment",     "Client must define the target business segment before Day 1 — scope ambiguity inflates cost and dilutes accuracy"),
    ("Benchmark Documents",        "Providing industry VSM benchmarks and transformation case studies pre-populates the RAG pipeline for higher-quality Phase 3 output"),
]
tb(sl, "Our Responsibility (Consulting Team)", 0.40, CT, 6.0, 0.28, sz=10, bold=True, color=BLUE)
tb(sl, "Client Responsibility", 6.75, CT, 6.0, 0.28, sz=10, bold=True, color=AMBER)
for i, (title, desc) in enumerate(csfs_ours):
    y = CT + 0.35 + i * 0.92
    rect(sl, 0.40, y, 6.20, 0.84, fill=LGRAY, line=BLUE, lw=Pt(0.6))
    rect(sl, 0.40, y, 0.14, 0.84, fill=BLUE)
    tb(sl, title, 0.62, y + 0.07, 5.82, 0.24, sz=10, bold=True, color=BLUE)
    tb(sl, desc,  0.62, y + 0.32, 5.82, 0.48, sz=10, color=DARK)
for i, (title, desc) in enumerate(csfs_client):
    y = CT + 0.35 + i * 0.76
    rect(sl, 6.75, y, 6.20, 0.68, fill=LGRAY, line=AMBER, lw=Pt(0.6))
    rect(sl, 6.75, y, 0.14, 0.68, fill=AMBER)
    tb(sl, title, 6.97, y + 0.06, 5.78, 0.24, sz=10, bold=True, color=AMBER)
    tb(sl, desc,  6.97, y + 0.30, 5.78, 0.34, sz=10, color=DARK)

tp(sl, [
    "Source system access (client CSF #2) is the single biggest engagement risk — make it a contractual requirement at SOW stage",
    "Domain SME availability is typically under-estimated: 2–3 hrs/week sounds small but clients routinely fail to resource it",
    "RAG quality management is our internal CSF most often overlooked: insist on benchmark documents before agents run in Phase 3",
    "C-suite sponsorship should be validated in the pre-sales stage — an engagement without it will stall at Phase 2",
])

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 18 — GTM Strategy
# ══════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Go-to-Market Strategy", "How we take this service offering to market, build a pipeline and establish market leadership", NAVY)
ftr(sl)

# Land & Expand model
tb(sl, "LAND & EXPAND MODEL", 0.40, CT, 4.0, 0.26, sz=10, bold=True, color=NAVY)
stages = [
    (RED,    "LAND",    "POC Engagement\n$45K, 2 weeks\n1 business unit\n1 org"),
    (AMBER,  "GROW",   "Full Engagement\n$280K, 8 weeks\nAll segments\n3–5 orgs"),
    (BLUE,   "EXPAND", "Enterprise License\n$500K+/year\nMulti-org platform\nContinuous AI"),
    (GREEN,  "SCALE",  "Platform-as-Service\n$150K/org/year\nSelf-serve teams\nMarketplace"),
]
for i, (col, stage, desc) in enumerate(stages):
    x = 0.40 + i * 3.10
    rect(sl, x, CT + 0.30, 2.90, 1.60, fill=col)
    tb(sl, stage, x, CT + 0.36, 2.90, 0.36, sz=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(sl, desc,  x + 0.14, CT + 0.76, 2.62, 1.08, sz=10, color=WHITE)
    if i < 3:
        tb(sl, "→", x + 2.90, CT + 0.90, 0.20, 0.32, sz=16, bold=True, color=MID, align=PP_ALIGN.CENTER)

# GTM Channels
tb(sl, "CHANNELS & TACTICS", 0.40, CT + 2.02, 4.0, 0.26, sz=10, bold=True, color=NAVY)
channels = [
    (BLUE,   "Direct Enterprise Sales", ["Target: CIO, CDO, Chief Transformation Officer","Entry via Digital Transformation or IT Strategy budget","$280K–$500K deal size","Cycle: 3–6 months"]),
    (PURPLE, "SI & Consulting Partnerships", ["Partner with Tier-2 SIs (not big-4 — they compete)","Co-sell model: SI brings client, we provide platform","Revenue share: 25–35% platform fee","Scale without headcount"]),
    (AMBER,  "Industry Associations", ["Banking: AFCA, ABA, FinTech Australia","Healthcare: HIMSS, ACHA","Publish benchmark reports for credibility","Speaking at transformation conferences"]),
    (GREEN,  "Thought Leadership / Content", ["'State of Digital Transformation Accuracy' annual report","Open benchmark database by industry","LinkedIn/GitHub open-source tools","Analyst briefings (Gartner, Forrester)"]),
]
for i, (col, name, pts) in enumerate(channels):
    x = 0.40 + i * 3.10
    rect(sl, x, CT + 2.30, 2.90, 2.30, fill=LGRAY, line=col, lw=Pt(0.8))
    rect(sl, x, CT + 2.30, 2.90, 0.30, fill=col)
    tb(sl, name, x + 0.10, CT + 2.34, 2.70, 0.22, sz=10, bold=True, color=WHITE)
    for j, p in enumerate(pts):
        oval(sl, x + 0.12, CT + 2.72 + j * 0.44, 0.08, 0.08, fill=col)
        tb(sl, p, x + 0.26, CT + 2.66 + j * 0.44, 2.56, 0.40, sz=10, color=DARK)

tp(sl, [
    "Always lead with a 2-week POC ($45K) — low enough to avoid procurement, high enough to show seriousness",
    "The POC should focus on ONE business unit the client cares about most — maximum impact, minimum scope",
    "SI partnerships are the scale engine: 1 SI partner can bring 10× the pipeline of a direct sales rep",
    "The 'State of DT Accuracy' annual report positions us as the thought leader in transformation intelligence — unique positioning",
])

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Target Industries (Value Props)
# ══════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Industry-Specific Value Propositions", "Tailored messaging and use cases for the five highest-value target verticals", NAVY)
ftr(sl)

industries = [
    (BLUE,   "🏦 Banking & Financial Services",
     ["SOX, FDIC, GLBA, PCI-DSS compliance mapping built-in","Digital banking capability gap analysis","Payments transformation value stream","Core banking modernisation roadmap","DORA Article 11 digital resilience compliance"],
     "Typical Deal Size: $350K–$600K  ·  Time-to-value: 4 weeks"),
    (GREEN,  "🏥 Healthcare & Life Sciences",
     ["HIPAA compliance framework pre-configured","EHR/EMR integration capability mapping","Clinical workflow value stream analysis","Patient journey digitisation roadmap","FHIR API discovery and mapping"],
     "Typical Deal Size: $280K–$500K  ·  Time-to-value: 6 weeks"),
    (PURPLE, "🛡️ Insurance",
     ["Claims processing VSM and bottleneck identification","Underwriting AI readiness assessment","APRA CPS 230 operational resilience mapping","Distribution channel capability inventory","RegTech compliance framework mapping"],
     "Typical Deal Size: $250K–$450K  ·  Time-to-value: 5 weeks"),
    (AMBER,  "🛒 Retail & eCommerce",
     ["Omnichannel value stream gap analysis","Supply chain digital capability mapping","Customer experience transformation roadmap","Loyalty platform capability discovery","GDPR data flow and compliance mapping"],
     "Typical Deal Size: $200K–$380K  ·  Time-to-value: 4 weeks"),
    (CYAN,   "📡 Telco & Media",
     ["BSS/OSS capability inventory and modernisation","5G product portfolio value stream mapping","Network API capability discovery","Content platform transformation roadmap","OSS/BSS consolidation business case"],
     "Typical Deal Size: $300K–$550K  ·  Time-to-value: 5 weeks"),
]
for i, (col, name, pts, note) in enumerate(industries):
    col_i = i % 3 if i < 3 else (i - 3)
    row_i = 0 if i < 3 else 1
    if i < 3:
        x = 0.40 + i * 4.22; y = CT
        w = 4.02; h = 2.06
    else:
        x = 0.40 + (i-3) * 6.38; y = CT + 2.14
        w = 6.18; h = 2.40
    rect(sl, x, y, w, h, fill=LGRAY, line=col, lw=Pt(0.8))
    rect(sl, x, y, w, 0.36, fill=col)
    tb(sl, name, x + 0.12, y + 0.06, w - 0.20, 0.26, sz=10, bold=True, color=WHITE)
    for j, pt in enumerate(pts):
        max_j = 4 if i >= 3 else 4
        if j <= max_j:
            oval(sl, x + 0.14, y + 0.48 + j * 0.32, 0.08, 0.08, fill=col)
            tb(sl, pt, x + 0.28, y + 0.42 + j * 0.32, w - 0.40, 0.28, sz=10, color=DARK)
    tb(sl, note, x + 0.14, y + h - 0.26, w - 0.20, 0.22, sz=10, bold=True, color=col, italic=True)

tp(sl, [
    "Banking is the highest-value vertical: regulatory burden creates budget urgency, and they have the largest IT estates",
    "Healthcare is fastest-growing vertical post-2023: EHR modernisation + FHIR API mandates create an immediate need",
    "Lead with industry-specific compliance frameworks in every pitch — generic transformation talk loses the room fast",
    "Build industry-specific reference architectures for each vertical — clients want to see 'we've done this in banking before'",
])

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Training for Consultants
# ══════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Consultant Training Programme", "Upskilling existing technology consultants to deliver the TransformHub Intelligence Service", NAVY)
ftr(sl)

modules = [
    (NAVY,   "Module 1",  "Platform Operations",    "2 days",
     ["Navigate all 8 modules as a power user","Configure organizations and segments","Run all 18 AI agents end-to-end","Interpret accuracy scores and module outputs","Seed and manage Context Hub documents"]),
    (BLUE,   "Module 2",  "AI Agent Configuration", "1 day",
     ["Understand LangGraph agent architecture","Configure RAG pipeline and document categories","Tune agent parameters for industry contexts","Diagnose and resolve agent output quality issues","Run BM25 and vector hybrid search"]),
    (GREEN,  "Module 3",  "Discovery Methodology",  "2 days",
     ["Conduct stakeholder interviews alongside AI","Validate AI-discovered products and capabilities","Facilitate human-in-the-loop review gates","Apply source triangulation principles","Write the D1.6 Executive Summary"]),
    (AMBER,  "Module 4",  "VSM Facilitation",       "1 day",
     ["Interpret L1/L2/L3 VSM outputs","Facilitate bottleneck workshops using AI data","Read Mermaid diagrams for executive audiences","Benchmark current-state FE% against industry data","Present VSM insights to non-technical stakeholders"]),
    (PURPLE, "Module 5",  "Executive Storytelling", "1 day",
     ["Structure the 5-phase engagement narrative","Build board-ready presentations from AI outputs","Handle 'black-box AI' objections with confidence","Translate accuracy scores into business language","Conduct a mock CXO debrief and Q&A"]),
]
tb(sl, "Training Duration: 7 days total  ·  Format: Blended (3 days instructor-led, 4 days hands-on platform lab)  ·  Certification: TransformHub Certified Practitioner (TCP)",
   0.40, CT, 12.53, 0.30, sz=10, bold=True, color=NAVY)
for i, (col, mod, title, dur, pts) in enumerate(modules):
    x = 0.40 + i * 2.52
    rect(sl, x, CT + 0.36, 2.38, 0.46, fill=col)
    tb(sl, mod, x + 0.10, CT + 0.40, 1.50, 0.20, sz=10, bold=True, color=WHITE)
    tb(sl, dur, x + 1.60, CT + 0.40, 0.70, 0.20, sz=10, color=WHITE, align=PP_ALIGN.RIGHT)
    tb(sl, title, x + 0.10, CT + 0.88, 2.18, 0.22, sz=10, bold=True, color=col)
    rect(sl, x, CT + 0.84, 2.38, 3.06, fill=LGRAY, line=col, lw=Pt(0.5))
    for j, pt in enumerate(pts):
        oval(sl, x + 0.12, CT + 1.12 + j * 0.52, 0.08, 0.08, fill=col)
        tb(sl, pt, x + 0.26, CT + 1.06 + j * 0.52, 2.05, 0.48, sz=10, color=DARK)

# Certification path
rect(sl, 0.40, CT + 3.98, 12.53, 0.60, fill=NAVY)
tb(sl, "📜  CERTIFICATION PATH:  Associate (after Module 1–3)  →  Practitioner (all 5 modules)  →  Lead (2 completed engagements)  →  Master Trainer (train others)",
   0.60, CT + 4.10, 12.10, 0.34, sz=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

tp(sl, [
    "Module 3 (Discovery Methodology) is the most important: consultants who can't interpret AI output will over-deliver raw data to clients",
    "Module 5 (Executive Storytelling) is the most under-invested in tech-first training programmes — make it mandatory",
    "Target: any existing technology or management consultant can achieve Associate certification in 3 days",
    "The Lead certification (after 2 engagements) ensures quality — only Leads should be assigned as Engagement Manager",
])

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 21 — Team Composition
# ══════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Team Composition", "Recommended team structure for a standard 8-week engagement — scalable up or down by scope", NAVY)
ftr(sl)

roles = [
    (NAVY,   "Engagement Manager",     "1 person",  "25%",  "$250/hr",  "$10,000",   ["Manages client relationship and gate reviews","Ensures delivery quality and timeline","Escalation point for access/scope issues","Presents at executive checkpoints"]),
    (BLUE,   "Sr. DT Consultant",      "1 person",  "100%", "$200/hr",  "$64,000",   ["Leads all 5 phases end-to-end","Interprets AI output for business context","Writes executive deliverables","Owns the transformation narrative"]),
    (GREEN,  "AI/ML Specialist",       "1 person",  "50%",  "$220/hr",  "$35,200",   ["Configures AI agents and RAG pipeline","Ensures accuracy scoring is calibrated","Diagnoses agent output quality issues","Tunes BM25 + vector search for each client"]),
    (AMBER,  "Business Analyst (×2)",  "2 people",  "100%", "$150/hr",  "$96,000",   ["Conducts stakeholder interviews","Documents and validates AI discoveries","Builds deliverable artefacts","Facilitates client workshops"]),
    (PURPLE, "Solution Architect",     "1 person",  "50%",  "$220/hr",  "$35,200",   ["Designs future state architecture","Validates technical feasibility of roadmap","Reviews AI-generated tech stack recommendations","Produces D3.5 Architecture Blueprint"]),
    (RED,    "Change Manager",         "1 person",  "25%",  "$180/hr",  "$14,400",   ["Stakeholder impact assessment","Change readiness evaluation","Communications plan","Training needs assessment"]),
]
for i, (col, role, ppl, alloc, rate, cost, pts) in enumerate(roles):
    col_i = i % 3; row_i = i // 3
    x = 0.40 + col_i * 4.22; y = CT + row_i * 1.86
    rect(sl, x, y, 4.02, 1.76, fill=LGRAY, line=col, lw=Pt(0.8))
    rect(sl, x, y, 4.02, 0.36, fill=col)
    tb(sl, role, x + 0.12, y + 0.06, 2.60, 0.24, sz=10, bold=True, color=WHITE)
    tb(sl, f"{ppl} · {alloc} · {rate}", x + 2.72, y + 0.10, 1.20, 0.20, sz=10, color=WHITE, align=PP_ALIGN.RIGHT)
    tb(sl, cost + " est.", x + 0.12, y + 0.42, 3.78, 0.22, sz=10, bold=True, color=col)
    for j, pt in enumerate(pts[:3]):
        oval(sl, x + 0.12, y + 0.72 + j * 0.34, 0.07, 0.07, fill=col)
        tb(sl, pt, x + 0.25, y + 0.66 + j * 0.34, 3.66, 0.30, sz=10, color=DARK)

# Total bar
rect(sl, 0.40, CT + 3.78, 12.53, 0.60, fill=NAVY)
tb(sl, "TOTAL:  4.5 FTE equivalent  ·  Est. consulting fees: $254,800  ·  Overhead & margin (35%): +$89,200  ·  TOTAL ENGAGEMENT: ~$344,000",
   0.55, CT + 3.92, 12.10, 0.28, sz=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

tp(sl, [
    "4.5 FTE is the minimum for an 8-week engagement — below this, quality degrades and consultant burnout increases",
    "The AI/ML Specialist at 50% allocation is often the most under-staffed role — consider 75% for first engagements",
    "Two Business Analysts are non-negotiable: one interfaces with AI output, one manages client communications",
    "Change Manager at 25% ensures the platform handover succeeds — clients without change support revert to old tools",
])

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 22 — Implementation Timeline
# ══════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Implementation Timeline", "Standard 8-week engagement schedule — with POC (2-week) and Extended (12-week) variants", NAVY)
ftr(sl)

weeks = [
    ("Week 1",  RED,    "DISCOVER", [("Org onboarding + system access","✓"),("AI Pass 1 — product discovery","✓"),("Stakeholder interviews (Day 3)","✓"),("Gate 1 review — product approval","✓")]),
    ("Week 2",  RED,    "DISCOVER", [("AI Pass 2 — capability mapping","✓"),("AI Pass 3 — functionality extraction","✓"),("Source attribution report","✓"),("D1.1–D1.6 Discovery package","✓")]),
    ("Week 3",  AMBER,  "MAP",      [("L1 segment VSM analysis","✓"),("L2 product-level VSM per capability","✓"),("Mermaid diagram generation (all)","✓"),("Benchmark upload to Context Hub","✓")]),
    ("Week 4",  AMBER,  "MAP",      [("L3 step classification analysis","✓"),("Bottleneck register","✓"),("Gate 2 review — VSM sign-off","✓"),("D2.1–D2.6 VSM package","✓")]),
    ("Week 5",  BLUE,   "DESIGN",   [("Future State Vision agent run","✓"),("Automation mix analysis","✓"),("3-band projection modelling","✓"),("RICE prioritisation scoring","✓")]),
    ("Week 6",  BLUE,   "DESIGN",   [("Architecture blueprint draft","✓"),("Gate 3 review — design sign-off","✓"),("D3.1–D3.6 Design package","✓"),("Exec presentation draft","✓")]),
    ("Week 7",  PURPLE, "GOVERN",   [("Risk identification — AI-assisted","✓"),("Compliance framework mapping","✓"),("SHA-256 audit trail setup","✓"),("D4.1–D4.6 Governance package","✓")]),
    ("Week 8",  GREEN,  "EXECUTE",  [("Transformation roadmap final","✓"),("Platform handover + training","✓"),("Executive presentation delivery","✓"),("D5.1–D5.6 Execution package","✓")]),
]
tb(sl, "STANDARD 8-WEEK ENGAGEMENT", 0.40, CT, 3.0, 0.26, sz=10, bold=True, color=NAVY)

# POC and Extended flags
rect(sl, 9.00, CT, 4.33, 0.26, fill=REDBG, line=RED, lw=Pt(0.5))
tb(sl, "⚡ POC Variant: Weeks 1–2 only ($45K)  ·  Extended: Add 4 weeks for complex enterprise ($180K+)",
   9.10, CT + 0.04, 4.10, 0.18, sz=10, color=RED, bold=True)

for i, (wk, col, phase, tasks) in enumerate(weeks):
    x = 0.40 + i * 1.60
    rect(sl, x, CT + 0.32, 1.50, 0.36, fill=col)
    tb(sl, wk, x + 0.08, CT + 0.36, 0.80, 0.22, sz=10, bold=True, color=WHITE)
    tb(sl, phase, x + 0.08, CT + 0.58, 1.34, 0.16, sz=6.5, color=WHITE)
    rect(sl, x, CT + 0.70, 1.50, 3.86, fill=LGRAY, line=col, lw=Pt(0.4))
    for j, (task, done) in enumerate(tasks):
        oval(sl, x + 0.10, CT + 0.84 + j * 0.82, 0.08, 0.08, fill=col)
        tb(sl, task, x + 0.24, CT + 0.78 + j * 0.82, 1.18, 0.36, sz=6.8, color=DARK)
        tb(sl, done, x + 1.26, CT + 0.84 + j * 0.82, 0.22, 0.22, sz=10, color=B4, bold=True)

# Gate markers
for i, (wk_num, label) in enumerate([(1, "Gate 1\nProducts"), (3, "Gate 2\nVSM"), (5, "Gate 3\nDesign"), (6, "Gate 4\nRisk")]):
    x = 0.40 + (wk_num * 2 - 0.5) * 0.80
    rect(sl, 0.40 + wk_num * 1.60 + 1.50, CT + 0.80, 0.14, 3.72, fill=MGRAY)
    tb(sl, label, 0.40 + wk_num * 1.60 + 1.56, CT + 1.80, 1.0, 0.40,
       sz=6.5, color=NAVY, align=PP_ALIGN.CENTER, bold=True)

tp(sl, [
    "Gate reviews are fixed — client must approve before next phase starts. This protects the engagement from scope drift",
    "The POC (Weeks 1–2 only, $45K) is the entry-point product: show instant value before asking for full engagement",
    "Weeks 7–8 can be parallelised for experienced teams: risk and roadmap activities overlap in practice",
    "Extended 12-week variant is recommended for orgs with 3+ business segments — add 4 weeks for each additional segment",
])

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 23 — Consulting Cost Model
# ══════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Consulting Investment Model", "Transparent pricing structure — three engagement tiers from POC to enterprise", NAVY)
ftr(sl)

tb(sl, "ENGAGEMENT TIERS", 0.40, CT, 3.0, 0.26, sz=10, bold=True, color=NAVY)

tiers = [
    (RED,    "⚡  POC Sprint",  "$38K – $52K",  "2 weeks",  "2.5 FTE",
     ["1 business unit, 1 org segment","Discovery Phase only (D1.1–D1.6)","Proof-of-value executive readout","Platform access for 30 days","Go/no-go decision support"]),
    (BLUE,   "⭐  Standard Engagement", "$280K – $360K", "8 weeks", "4.5 FTE",
     ["All 5 phases, all deliverables","Full org discovery (all segments)","Risk & compliance mapping included","Platform licence (12 months)","2-day client team training"]),
    (GREEN,  "🏢  Enterprise Programme", "$480K – $650K", "12–16 weeks", "6 FTE",
     ["3–5 business units or orgs","Custom compliance framework","Architecture review board sessions","Ongoing AI agent monitoring","Quarterly accuracy reviews"]),
]
for i, (col, name, price, dur, fte, items) in enumerate(tiers):
    x = 0.40 + i * 4.22
    rect(sl, x, CT + 0.30, 4.02, 0.48, fill=col)
    tb(sl, name,  x + 0.14, CT + 0.34, 3.74, 0.26, sz=10, bold=True, color=WHITE)
    rect(sl, x, CT + 0.78, 4.02, 3.10, fill=LGRAY, line=col, lw=Pt(0.8))
    tb(sl, price, x, CT + 0.82, 4.02, 0.46, sz=22, bold=True, color=col, align=PP_ALIGN.CENTER)
    tb(sl, f"{dur}  ·  {fte} equivalent", x, CT + 1.28, 4.02, 0.24, sz=10, color=MID, align=PP_ALIGN.CENTER)
    for j, it in enumerate(items):
        oval(sl, x + 0.16, CT + 1.62 + j * 0.44, 0.08, 0.08, fill=col)
        tb(sl, it, x + 0.30, CT + 1.56 + j * 0.44, 3.60, 0.40, sz=10, color=DARK)

# Cost breakdown row
tb(sl, "INDICATIVE COST BREAKDOWN — STANDARD ENGAGEMENT ($320K mid-point)", 0.40, CT + 3.98, 8.0, 0.26, sz=10, bold=True, color=NAVY)
breakdown = [
    ("Engagement Manager",  "$10K",   "3.1%",  NAVY),
    ("Sr. DT Consultant",   "$64K",   "20.0%", BLUE),
    ("AI/ML Specialist",    "$35K",   "10.9%", GREEN),
    ("Business Analysts ×2","$96K",   "30.0%", AMBER),
    ("Solution Architect",  "$35K",   "10.9%", PURPLE),
    ("Change Manager",      "$14K",   "4.4%",  RED),
    ("Overhead & Margin",   "$66K",   "20.6%", MID),
]
rect(sl, 0.40, CT + 4.26, 12.53, 0.36, fill=LGRAY, line=MGRAY, lw=Pt(0.5))
for i, (role, cost, pct, col) in enumerate(breakdown):
    x = 0.50 + i * 1.78
    tb(sl, cost, x, CT + 4.30, 1.60, 0.20, sz=10, bold=True, color=col, align=PP_ALIGN.CENTER)
    tb(sl, pct,  x, CT + 4.48, 1.60, 0.14, sz=10, color=MID, align=PP_ALIGN.CENTER)

tp(sl, [
    "Lead with the POC ($38–52K) — most procurement teams can approve this without a full committee review",
    "Standard engagement at $280–360K delivers 3–5× ROI in 18 months — present this as an investment, not a cost",
    "Overhead and margin at 35% is conservative for technology consulting — adjust per firm's margin targets",
    "Consider a 'Platform + Consulting' bundle: platform licence ($80K) + consulting ($280K) priced together at $320K saves the client admin",
])

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 24 — Client Technology Cost
# ══════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Client Technology Investment", "Full cost picture for the client — platform, licences, infrastructure and internal resources", NAVY)
ftr(sl)

tb(sl, "TECHNOLOGY & PLATFORM COSTS", 0.40, CT, 5.0, 0.26, sz=10, bold=True, color=NAVY)
tech = [
    ("TransformHub Platform Licence",  "Annual",      "$80K – $150K",   BLUE,
     "Per-org licence including all 18 AI agents, unlimited users, full feature set. Volume discounts for 3+ orgs."),
    ("OpenAI API (GPT-4o)",            "Per engmt.",  "$2K – $8K",      GREEN,
     "Varies by number of orgs, runs and context doc volume. ~$2K for a standard 8-week engagement."),
    ("Cloud Hosting (AWS/Azure/GCP)",  "Monthly",     "$500 – $2,000",  BLUE,
     "Next.js on Vercel/App Service + FastAPI on Fargate/Container Apps + PostgreSQL managed DB."),
    ("PostgreSQL + pgvector",          "Included",    "$0",             GREEN,
     "Open-source database — included in cloud hosting estimate above. No additional licence required."),
    ("Prisma ORM",                     "Included",    "$0",             GREEN,
     "MIT licence — no cost. Enterprise support available at $5K/year if required."),
]
for i, (item, freq, cost, col, note) in enumerate(tech):
    y = CT + 0.30 + i * 0.64
    rect(sl, 0.40, y, 12.53, 0.60, fill=(LGRAY if i % 2 == 0 else WHITE), line=MGRAY, lw=Pt(0.4))
    tb(sl, item, 0.55, y + 0.18, 4.80, 0.26, sz=10, bold=True, color=DARK)
    tb(sl, freq, 5.50, y + 0.18, 1.40, 0.26, sz=10, color=MID, align=PP_ALIGN.CENTER)
    tb(sl, cost, 7.00, y + 0.16, 2.20, 0.28, sz=10, bold=True, color=col, align=PP_ALIGN.CENTER)
    tb(sl, note, 9.30, y + 0.18, 3.60, 0.26, sz=10, color=MID, italic=True)

# Client internal costs
tb(sl, "CLIENT INTERNAL COSTS (Indicative)", 0.40, CT + 3.56, 5.0, 0.26, sz=10, bold=True, color=NAVY)
for i, (item, cost, note) in enumerate([
    ("Stakeholder time (10 ppl × 3hrs/week × 8 weeks)", "$30K – $50K",  "Based on avg $125/hr fully loaded internal cost"),
    ("IT environment setup (VPN, access, CI/CD)",        "$10K – $20K",  "One-off infrastructure preparation"),
    ("Internal training (client team, 5 people)",        "$5K – $15K",   "4 hours each + internal facilitation time"),
]):
    y = CT + 3.84 + i * 0.36
    rect(sl, 0.40, y, 12.53, 0.32, fill=(LGRAY if i % 2 == 0 else WHITE))
    tb(sl, item, 0.55, y + 0.06, 7.50, 0.22, sz=10, color=DARK)
    tb(sl, cost, 8.20, y + 0.06, 2.20, 0.22, sz=10, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
    tb(sl, note, 10.50, y + 0.06, 2.30, 0.22, sz=10, color=MID, italic=True)

# Total investment
rect(sl, 0.40, CT + 4.94, 12.53, 0.54, fill=NAVY)
tb(sl, "TOTAL CLIENT INVESTMENT (Standard Engagement):  Technology ~$100K–$180K  ·  Internal ~$45K–$85K  ·  Consulting ~$320K  ·  TOTAL: ~$465K–$585K",
   0.55, CT + 5.08, 12.10, 0.26, sz=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

tp(sl, [
    "The platform licence ($80–150K) is the largest technology cost — position it as replacing 6 tools at $40–80K each = cost-neutral",
    "OpenAI API cost of $2–8K is surprisingly low: emphasise this to CFOs who expect AI to be expensive",
    "Total client investment of $465–585K over 18 months delivers $7.6M+ in quantified value — present the ROI slide next",
    "Internal cost of $45–85K is often invisible to procurement — surfacing it creates a sense of 'total commitment' that builds buy-in",
])

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 25 — Business Case & ROI
# ══════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Client Business Case & ROI", "Typical value delivered in 18 months — based on pilot engagements and industry benchmarks", NAVY)
ftr(sl)

tb(sl, "VALUE DELIVERED", 0.40, CT, 3.0, 0.26, sz=10, bold=True, color=B4)
value_items = [
    ("Reduced Discovery Cost",    "$250K – $400K",  "vs traditional 12-week consulting workshop approach at $500K+ typical",             GREEN),
    ("Faster Time-to-Market",     "$500K – $2M",    "6 months earlier delivery × $83K–$333K/month digital revenue per product",          GREEN),
    ("Risk Avoidance",            "$500K – $5M",    "2 CRITICAL risks mitigated pre-transformation × $250K–$2.5M remediation cost each", GREEN),
    ("Tool Consolidation Savings","$200K – $400K/yr","6 tools × $33–67K/year replaced by one platform (ongoing annual saving)",          GREEN),
    ("Analyst Productivity Gain", "$150K – $300K",  "8× productivity improvement × 2 analysts × $75K–$150K/year fully-loaded cost",      GREEN),
]
rect(sl, 0.40, CT + 0.30, 12.53, 0.34, fill=NAVY)
for j, h in enumerate(["Value Driver", "Estimated Value", "Basis of Estimate"]):
    aligns = [PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.LEFT]
    tb(sl, h, [0.55, 5.20, 7.80][j], CT + 0.36, [4.50, 2.40, 4.80][j], 0.22,
       sz=10, bold=True, color=WHITE, align=aligns[j])
for i, (item, val, basis, col) in enumerate(value_items):
    y = CT + 0.66 + i * 0.52
    rect(sl, 0.40, y, 12.53, 0.50, fill=(LGRAY if i % 2 == 0 else WHITE))
    tb(sl, item, 0.55, y + 0.14, 4.50, 0.24, sz=10, bold=True, color=DARK)
    tb(sl, val,  5.20, y + 0.12, 2.40, 0.28, sz=10, bold=True, color=col, align=PP_ALIGN.CENTER)
    tb(sl, basis, 7.80, y + 0.14, 4.70, 0.24, sz=10, color=MID, italic=True)

# Total value and cost boxes
rect(sl, 0.40, CT + 3.32, 5.90, 1.26, fill=GREENBG, line=GREEN, lw=Pt(1.2))
tb(sl, "TOTAL 18-MONTH VALUE", 0.60, CT + 3.40, 5.50, 0.26, sz=10, bold=True, color=B4)
tb(sl, "$1.6M – $8.1M", 0.60, CT + 3.68, 5.50, 0.56, sz=32, bold=True, color=B4, align=PP_ALIGN.CENTER)
tb(sl, "Conservative scenario shown. Expected: $5.2M  ·  Optimistic: $8.1M",
   0.60, CT + 4.24, 5.50, 0.26, sz=10, color=MID, align=PP_ALIGN.CENTER, italic=True)

rect(sl, 6.55, CT + 3.32, 2.88, 1.26, fill=REDBG, line=RED, lw=Pt(1.2))
tb(sl, "TOTAL INVESTMENT", 6.70, CT + 3.40, 2.60, 0.26, sz=10, bold=True, color=RED)
tb(sl, "~$500K", 6.70, CT + 3.68, 2.60, 0.56, sz=28, bold=True, color=RED, align=PP_ALIGN.CENTER)
tb(sl, "Consulting + Tech + Internal",
   6.70, CT + 4.24, 2.60, 0.26, sz=10, color=MID, align=PP_ALIGN.CENTER, italic=True)

rect(sl, 9.60, CT + 3.32, 3.33, 1.26, fill=BLUEBG, line=BLUE, lw=Pt(1.5))
tb(sl, "ROI  (18 months)", 9.78, CT + 3.40, 2.96, 0.26, sz=10, bold=True, color=BLUE)
tb(sl, "3× – 15×", 9.78, CT + 3.68, 2.96, 0.56, sz=28, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
tb(sl, "Payback period: 4–7 months",
   9.78, CT + 4.24, 2.96, 0.26, sz=10, color=MID, align=PP_ALIGN.CENTER, italic=True)

tp(sl, [
    "Lead with the 'faster time-to-market' row — it's the largest value driver and easiest for business leaders to internalise",
    "Risk avoidance is the most credible number for regulated industries — one prevented compliance failure justifies the entire investment",
    "The 4–7 month payback period is the CFO close: 'You recover your investment before the engagement is finished'",
    "Always customise this slide with client-specific numbers during pre-sales — generic ROI slides are ignored, specific ones are presented to boards",
])

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 26 — Call to Action
# ══════════════════════════════════════════════════════════════════════════
sl = ns()
rect(sl, 0, 0, 13.33, 3.50, fill=NAVY)
rect(sl, 0, 3.50, 13.33, 0.08, fill=AMBER)

tb(sl, "Ready to Transform with Intelligence?", 0.70, 0.40, 11.93, 0.80,
   sz=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(sl, "Three ways to get started — choose the right entry point for your client",
   1.0, 1.30, 11.33, 0.36, sz=14, color=B8,
   align=PP_ALIGN.CENTER, italic=True)

options = [
    (RED,    "⚡",  "2-Week POC Sprint",    "$38K – $52K",
     ["1 business unit, 1 segment","Full discovery + D1 package","Executive readout + demo","Platform access 30 days","Go/no-go recommendation"]),
    (BLUE,   "⭐",  "8-Week Full Engagement","$280K – $360K",
     ["All 5 phases, 25+ deliverables","Full org + compliance mapping","Board-ready executive deck","12-month platform licence","Team training included"]),
    (GREEN,  "🏢",  "Enterprise Programme", "$480K – $650K",
     ["3–5 orgs or business units","Custom compliance frameworks","Architecture review sessions","Ongoing AI monitoring","Quarterly accuracy reviews"]),
]
for i, (col, icon, name, price, pts) in enumerate(options):
    x = 0.55 + i * 4.12
    rect(sl, x, 3.72, 3.92, 3.28, fill=LGRAY, line=col, lw=Pt(1.2))
    rect(sl, x, 3.72, 3.92, 0.44, fill=col)
    tb(sl, f"{icon}  {name}", x + 0.14, 3.76, 3.64, 0.36, sz=10, bold=True, color=WHITE)
    tb(sl, price, x, 4.18, 3.92, 0.44, sz=20, bold=True, color=col, align=PP_ALIGN.CENTER)
    for j, pt in enumerate(pts):
        oval(sl, x + 0.16, 4.66 + j * 0.46, 0.08, 0.08, fill=col)
        tb(sl, pt, x + 0.30, 4.60 + j * 0.46, 3.52, 0.42, sz=10, color=DARK)

ftr(sl, "TransformHub Intelligence Service  ·  AI-Powered Digital Transformation  ·  transformhub.ai")
tp(sl, [
    "Always close with: 'Which business unit would you most want to transform first?' — it moves from abstract to concrete immediately",
    "The POC is designed to be approved within 2 weeks — keep a SOW template ready to send within 24 hours of interest",
    "Pre-built POC scope: Banking = Digital Lending segment, Healthcare = Patient Engagement, Insurance = Claims Processing",
    "Follow-up email within 24 hours: attach this deck + a one-page POC scope document tailored to their industry",
])

# ══════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════
OUT = "/Users/125066/projects/TransformHub/docs/TransformHub_Service_Offering.pptx"
prs.save(OUT)
print(f"✅  Saved: {OUT}")
print(f"   {len(prs.slides)} slides")
